# app/services/report/report_formatters/report_pptx_formatter.py
from typing import Dict, Any, List, Optional, Tuple
from io import BytesIO
import logging
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

from ..report_formatter import ReportFormatter

logger = logging.getLogger(__name__)

class ReportPptxFormatter(ReportFormatter):
    """Formatter for PowerPoint (PPTX) reports"""
    
    def generate(self) -> BytesIO:
        """
        Generate a PPTX report
        
        Returns:
            BytesIO buffer with the PPTX report
        """
        prs = Presentation()
        buffer = BytesIO()
        
        # Constants for slide layout
        DEFAULT_CHART_TOP = Inches(1.5)
        DEFAULT_CHART_LEFT = Inches(1.0)
        DEFAULT_CHART_WIDTH = Inches(8.0)
        DEFAULT_CHART_HEIGHT = Inches(4.5)
        DEFAULT_TABLE_ROWS = 10
        
        # Find first valid report type to focus on
        primary_report = self.get_first_valid_report()
        
        if not primary_report:
            # Create error presentation
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = "Report Error"
            
            error_msg = self.get_all_errors()
            text_box = slide.shapes.add_textbox(
                Inches(1.0), Inches(2.0), Inches(8.0), Inches(4.0)
            )
            tf = text_box.text_frame
            p = tf.add_paragraph()
            p.text = error_msg
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(255, 0, 0)
            
            prs.save(buffer)
            buffer.seek(0)
            return buffer
        
        # Create title slide
        slide_title = prs.slides.add_slide(prs.slide_layouts[0])
        slide_title.shapes.title.text = self.report_title
        
        try:
            subtitle = slide_title.placeholders[1]
            entity_title = primary_report['result'].get('params', {}).get(
                "sheet_name", primary_report['type'].replace("_", " ").title()
            )
            subtitle.text = f"Entity: {entity_title}\nGenerated: {self.generation_timestamp}"
        except (IndexError, KeyError):
            logger.warning("Could not set subtitle in title slide")
        
        # Process report data
        report_type = primary_report['type']
        result = primary_report['result']
        analysis = result.get('analysis', {})
        params = result.get('params', {})
        data = result.get('data', [])
        
        # Add summary slide
        self._add_summary_slide(prs, report_type, result)
        
        # Add charts
        charts = analysis.get('charts', {})
        for chart_key, chart_bytes in charts.items():
            self._add_chart_slide(prs, chart_key, chart_bytes)
        
        # Add data table if requested
        if data and params.get("include_data_table_in_ppt", False):
            self._add_data_table_slide(prs, report_type, result)
            
        # Add insights slide
        if analysis.get('insights'):
            self._add_insights_slide(prs, analysis.get('insights', {}))
        
        # Add conclusion slide
        self._add_conclusion_slide(prs, report_type, result)
        
        # Save presentation
        try:
            prs.save(buffer)
            buffer.seek(0)
            return buffer
        except Exception as e:
            logger.error(f"Error saving PPTX: {e}")
            
            # Create error presentation
            buffer = BytesIO()
            error_prs = Presentation()
            slide = error_prs.slides.add_slide(error_prs.slide_layouts[5])
            slide.shapes.title.text = "Error Generating PPTX Report"
            
            text_box = slide.shapes.add_textbox(
                Inches(1.0), Inches(2.0), Inches(8.0), Inches(4.0)
            )
            tf = text_box.text_frame
            p = tf.add_paragraph()
            p.text = str(e)
            p.font.color.rgb = RGBColor(255, 0, 0)
            
            try:
                error_prs.save(buffer)
                buffer.seek(0)
                return buffer
            except Exception:
                return BytesIO(b"Failed to generate PPTX report.")
    
    def _add_summary_slide(self, prs: Presentation, report_type: str, result: Dict[str, Any]) -> None:
        """Add summary slide with key statistics"""
        analysis = result.get('analysis', {})
        params = result.get('params', {})
        stats = analysis.get('summary_stats', {})
        
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content layout
        entity_title = params.get("sheet_name", report_type.replace("_", " ").title())
        slide.shapes.title.text = f"{entity_title} - Summary"
        
        # Create content for summary
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        # Add record count
        if 'record_count' in stats:
            p = text_frame.add_paragraph()
            p.text = f"Total Records: {stats['record_count']}"
            p.font.bold = True
            p.level = 0
        
        # Add simple stats (non-dict values)
        simple_stats = {
            k: v for k, v in stats.items()
            if not k.startswith('_') and 
            k != 'record_count' and
            not isinstance(v, (dict, list)) and
            v is not None
        }
        
        if simple_stats:
            p = text_frame.add_paragraph()
            p.text = "Key Statistics:"
            p.font.bold = True
            p.level = 0
            
            for key, value in simple_stats.items():
                p = text_frame.add_paragraph()
                p.text = f"{key.replace('_', ' ').title()}: {value}"
                p.level = 1
    
    def _add_chart_slide(self, prs: Presentation, chart_key: str, chart_bytes: Optional[BytesIO]) -> None:
        """Add a slide with a chart"""
        if not isinstance(chart_bytes, BytesIO):
            logger.warning(f"Invalid chart data for '{chart_key}'")
            return
            
        chart_bytes.seek(0)
        chart_title = chart_key.replace('_', ' ').title()
        
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only layout
            slide.shapes.title.text = chart_title
            
            # Add chart image
            pic_left, pic_top = Inches(0.5), Inches(1.5)
            pic = slide.shapes.add_picture(chart_bytes, pic_left, pic_top)
            
            # Resize and center image
            pic_width, pic_height = pic.width, pic.height
            max_width, max_height = Inches(9.0), Inches(5.0)
            
            # Scale down if too large
            if pic_width > max_width:
                ratio = max_width / pic_width
                pic_width = max_width
                pic_height = int(pic_height * ratio)
                
            if pic_height > max_height:
                ratio = max_height / pic_height
                pic_height = max_height
                pic_width = int(pic_width * ratio)
                
            pic.width, pic.height = pic_width, pic_height
            
            # Center horizontally
            pic.left = int((prs.slide_width - pic_width) / 2)
            
        except Exception as e:
            logger.error(f"Error adding chart '{chart_key}' to PPTX: {e}")
    
    def _add_data_table_slide(self, prs: Presentation, report_type: str, result: Dict[str, Any]) -> None:
        """Add a slide with a data table"""
        params = result.get('params', {})
        data = result.get('data', [])
        columns = params.get('columns', [])
        
        if not data or not columns:
            return
            
        entity_title = params.get("sheet_name", report_type.replace("_", " ").title())
        max_rows = params.get("max_ppt_table_rows", 10)
        
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only layout
            slide.shapes.title.text = f"{entity_title} - Data Sample"
            
            # Select important columns to display (prioritize key fields)
            display_cols = []
            priority_cols = ["id", "name", "title", "submitted_by", "created_at", "status"]
            
            for col in priority_cols:
                if col in columns:
                    display_cols.append(col)
                    
            # If we don't have enough columns, add more
            if len(display_cols) < 3:
                # Add other non-ANSWERS_PREFIX columns until we have at least 3
                for col in columns:
                    if col not in display_cols and not col.startswith("answers."):
                        display_cols.append(col)
                        if len(display_cols) >= 5:  # Max columns for readability
                            break
            
            # Create table
            rows = min(len(data) + 1, max_rows + 1)  # +1 for header
            cols = len(display_cols)
            
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(9.0)
            height = Inches(5.0)
            
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table
            
            # Add header row
            for i, col in enumerate(display_cols):
                cell = table.cell(0, i)
                clean_name = col.replace('.', ' ').replace('_', ' ').title()
                cell.text = clean_name
                
                # Format header
                para = cell.text_frame.paragraphs[0]
                para.font.bold = True
                para.font.size = Pt(11)
                para.alignment = PP_ALIGN.CENTER
                
                # Set fill
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0, 112, 192)  # Blue
                
                # Set text color
                para.font.color.rgb = RGBColor(255, 255, 255)  # White
            
            # Add data rows
            for i, row_dict in enumerate(data[:max_rows]):
                for j, col in enumerate(display_cols):
                    cell = table.cell(i + 1, j)
                    val = row_dict.get(col)
                    
                    if isinstance(val, bool):
                        val = "Yes" if val else "No"
                    elif val is None:
                        val = ""
                    
                    # Truncate long values
                    str_val = str(val)
                    if len(str_val) > 50:
                        str_val = str_val[:47] + "..."
                        
                    cell.text = str_val
                    
                    # Format data cell
                    para = cell.text_frame.paragraphs[0]
                    para.font.size = Pt(10)
                    
                    # Alternate row colors
                    if i % 2 == 1:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(235, 241, 252)  # Light blue
            
            # Add note if data was truncated
            if len(data) > max_rows:
                note_top = table.cell(0, 0).top + table.height + Inches(0.1)
                if note_top < prs.slide_height - Inches(0.5):
                    txBox = slide.shapes.add_textbox(
                        left, note_top, width, Inches(0.5)
                    )
                    tf = txBox.text_frame
                    p = tf.add_paragraph()
                    p.text = f"Note: Showing {max_rows} of {len(data)} records."
                    p.font.size = Pt(10)
                    p.font.italic = True
                    p.alignment = PP_ALIGN.CENTER
                    
        except Exception as e:
            logger.error(f"Error adding data table to PPTX: {e}")
    
    def _add_insights_slide(self, prs: Presentation, insights: Dict[str, str]) -> None:
        """Add a slide with key insights"""
        if not insights:
            return
            
        # Skip status insight which is often just a placeholder
        cleaned_insights = {k: v for k, v in insights.items() if k != 'status'}
        if not cleaned_insights:
            return
            
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content layout
        slide.shapes.title.text = "Key Insights"
        
        # Add insights as bullet points
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        for key, insight_text in cleaned_insights.items():
            p = text_frame.add_paragraph()
            p.text = insight_text
            p.level = 0
            p.font.size = Pt(18)
    
    def _add_conclusion_slide(self, prs: Presentation, report_type: str, result: Dict[str, Any]) -> None:
        """Add a concluding slide with summary"""
        analysis = result.get('analysis', {})
        stats = analysis.get('summary_stats', {})
        insights = analysis.get('insights', {})
        
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content layout
        slide.shapes.title.text = "Conclusions"
        
        # Add conclusions as bullet points
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        # Add record count
        record_count = stats.get('record_count')
        entity_name = result.get('params', {}).get("sheet_name", report_type.replace("_", " ").title())
        
        p = text_frame.add_paragraph()
        if record_count is not None:
            p.text = f"Analysis completed for {record_count} {entity_name} records."
        else:
            p.text = f"Analysis completed for {entity_name}."
        p.font.bold = True
        p.level = 0
        
        # Add key insight for conclusion if available
        top_insight_keys = [
            'primary_finding', 'top_category_info', 'dominant_role', 
            'primary_environment', 'top_user', 'top_form', 'creator_activity'
        ]
        
        added_insights = 0
        
        # First try to find insights with priority keys
        for key in top_insight_keys:
            if key in insights and added_insights < 3:
                p = text_frame.add_paragraph()
                p.text = insights[key]
                p.level = 1
                added_insights += 1
        
        # If we don't have enough, add other insights
        if added_insights < 3:
            for key, text in insights.items():
                if key not in top_insight_keys and key != 'status' and added_insights < 3:
                    p = text_frame.add_paragraph()
                    p.text = text
                    p.level = 1
                    added_insights += 1
        
        # Add thank you
        p = text_frame.add_paragraph()
        p.text = "Thank you for reviewing this report."
        p.font.italic = True
        p.level = 0
        p.alignment = PP_ALIGN.CENTER