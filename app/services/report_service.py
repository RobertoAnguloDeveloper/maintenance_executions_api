# app/services/report_service.py
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
import matplotlib
matplotlib.use('Agg') # Use 'Agg' backend BEFORE importing pyplot
import xlsxwriter # Ensure this is imported
import csv
import zipfile
from io import BytesIO, StringIO
from typing import List, Dict, Tuple, Optional, Any, Type, Callable, Union
from datetime import date, datetime
import logging
import traceback
import os

# SQLAlchemy Imports
from sqlalchemy import inspect as sqla_inspect, or_, asc, desc, func, text, Boolean, String, Integer, Date, DateTime
from sqlalchemy.sql.schema import Column as AlchemyColumn
from sqlalchemy.orm import Query, joinedload, selectinload, Session, aliased
from app import db # Assuming db is initialized SQLAlchemy instance

# App-specific Imports
from app.models import (
    User, FormSubmission, AnswerSubmitted, Form, Role, Environment,
    Question, Answer, QuestionType, Permission, RolePermission,
    FormQuestion, FormAnswer, Attachment, ReportTemplate # Add any other models you might report on
)
from app.utils.permission_manager import PermissionManager, EntityType, RoleType

# ReportLab Imports
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as RLImage, Table, TableStyle, PageBreak, BaseDocTemplate
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
from reportlab.lib.pagesizes import letter, landscape
from reportlab.lib.units import inch
from reportlab.lib import colors

# DOCX Imports
import docx # Import the whole module
from docx.shared import Inches, Pt, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.document import Document as DocxDocumentClass

# PPTX Imports
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor as PptxRGBColor

logger = logging.getLogger(__name__)

# --- Type Aliases ---
DataList = List[Dict[str, Any]]
AnalysisDict = Dict[str, Any]
ParamsDict = Dict[str, Any]
ProcessedData = Dict[str, Dict[str, Any]]
ReportResult = Tuple[Optional[BytesIO], Optional[str], Optional[str], Optional[str]]
SchemaResult = Tuple[Optional[Dict[str, Any]], Optional[str]]
StatsGenerator = Callable[[pd.DataFrame, ParamsDict], Dict[str, Any]]
ChartGenerator = Callable[[pd.DataFrame, AnalysisDict], Dict[str, BytesIO]] # Takes analysis dict now
InsightGenerator = Callable[[pd.DataFrame, AnalysisDict, ParamsDict], Dict[str, str]]
PptxGenerator = Callable[[Dict[str, Any], ParamsDict], BytesIO]
DocxGenerator = Callable[[ProcessedData, ParamsDict], BytesIO]
PdfGenerator = Callable[[ProcessedData, ParamsDict], BytesIO]
XlsxGenerator = Callable[[ProcessedData, ParamsDict], BytesIO]
CsvGenerator = Callable[[ProcessedData, ParamsDict], BytesIO]

# --- Constants ---
SUPPORTED_FORMATS = ["xlsx", "csv", "pdf", "docx", "pptx"]
MULTI_ENTITY_FORMATS = ["xlsx", "csv", "pdf", "docx"]
VISUAL_FORMATS = ["pdf", "docx", "pptx"]
DEFAULT_REPORT_TITLE = "Data Analysis Report"
DEFAULT_CHART_WIDTH = 6.5 * inch
DEFAULT_CHART_HEIGHT = 3.5 * inch
DEFAULT_PPTX_CHART_WIDTH = PptxInches(8)
DEFAULT_PPTX_CHART_HEIGHT = PptxInches(4.5)
DEFAULT_PPTX_CHART_TOP = PptxInches(1.5)
DEFAULT_PPTX_CHART_LEFT = PptxInches(1)
DEFAULT_PPTX_TABLE_ROWS = 10
MAX_XLSX_SHEET_NAME_LEN = 31
ANSWERS_PREFIX = "answers."
GENERIC_CATEGORICAL_COLS = ['type', 'name', 'status', 'action', 'role', 'environment', 'is_public', 'is_deleted', 'is_super_user', 'is_signature', 'file_type', 'question_type']
MAX_UNIQUE_GENERIC_CHART = 15


class ReportService:
    """
    Service responsible for generating customizable reports for various entities,
    respecting user permissions. Supports XLSX, CSV, PDF, DOCX, PPTX formats.
    Handles multi-entity requests and leverages question types for analysis.
    Provides basic analysis for all entities in visual reports.
    Defaults to all direct table columns for reports when columns are not specified.
    Restricts sensitive columns (e.g., password_hash) for non-admin users.
    """

    ENTITY_CONFIG = {
        'form_submissions': {
            'model': FormSubmission, 'view_permission_entity': EntityType.SUBMISSIONS,
            'default_columns': [ "id", "form_id", "form.title", "submitted_by", "user.email", "user.first_name", "user.last_name", "submitted_at", "updated_at"], # Added user details
            'analysis_hints': {
                'date_columns': ['submitted_at', 'created_at', 'updated_at'],
                'categorical_columns': ['submitted_by', 'form.title', 'user.email'],  # Added user.email
                'dynamic_answer_prefix': ANSWERS_PREFIX,
            },
            'default_sort': [{"field": "submitted_at", "direction": "desc"}],
            'stats_generators': ['generate_submission_stats'], # Kept specific due to dynamic answers complexity
            'chart_generators': ['generate_submission_charts'],# Kept specific
            'insight_generators': ['generate_submission_insights'],# Kept specific
            'format_generators': { 'pptx': '_generate_submission_pptx' }
        },
        "users": {
            'model': User, 'view_permission_entity': EntityType.USERS,
            'default_columns': [ "id", "username", "first_name", "last_name", "email", "role.name", "environment.name", "created_at" ],
            'analysis_hints': {
                'date_columns': ['created_at', 'updated_at'],
                'categorical_columns': ['role.name', 'environment.name', 'is_deleted', 'username', 'email'], # Added username, email
            },
            'default_sort': [{"field": "username", "direction": "asc"}],
            'stats_generators': ['generate_user_stats', '_generate_generic_stats'],       # MODIFIED
            'chart_generators': ['generate_user_charts', '_generate_generic_charts'],     # MODIFIED
            'insight_generators': ['generate_user_insights', '_generate_generic_insights'], # MODIFIED
            'format_generators': {}
        },
        "forms": {
            'model': Form, 'view_permission_entity': EntityType.FORMS,
            'default_columns': [ "id", "title", "description", "creator.username", "creator.environment.name", "is_public", "created_at" ],
            'analysis_hints': {
                'date_columns': ['created_at', 'updated_at'],
                'categorical_columns': ['creator.username', 'creator.environment.name', 'is_public', 'is_deleted', 'title'], # Added title
            },
            'default_sort': [{"field": "title", "direction": "asc"}],
            'stats_generators': ['generate_form_stats', '_generate_generic_stats'],       # MODIFIED
            'chart_generators': ['generate_form_charts', '_generate_generic_charts'],     # MODIFIED
            'insight_generators': ['generate_form_insights', '_generate_generic_insights'], # MODIFIED
            'format_generators': {}
        },
        "environments": {
            'model': Environment, 'view_permission_entity': EntityType.ENVIRONMENTS,
            'default_columns': ["id", "name", "description", "created_at"],
            'analysis_hints': {
                'date_columns': ['created_at', 'updated_at'],
                'categorical_columns': ['name']
            },
            'default_sort': [{"field": "name", "direction": "asc"}],
            'stats_generators': ['generate_environment_stats', '_generate_generic_stats'],
            'chart_generators': ['_generate_generic_charts'], # Already generic focused
            'insight_generators': ['_generate_generic_insights'],
            'format_generators': {}
        },
        "roles": {
            'model': Role, 'view_permission_entity': EntityType.ROLES,
            'default_columns': ["id", "name", "description", "is_super_user", "created_at"],
            'analysis_hints': {
                'categorical_columns': ['is_super_user', 'name']
            },
            'default_sort': [{"field": "name", "direction": "asc"}],
            'stats_generators': ['generate_role_stats', '_generate_generic_stats'],
            'chart_generators': ['_generate_generic_charts'], # Already generic focused
            'insight_generators': ['_generate_generic_insights'],
            'format_generators': {}
        },
        "answers_submitted": { # This refers to the direct model AnswerSubmitted
            'model': AnswerSubmitted, 'view_permission_entity': EntityType.SUBMISSIONS, # Permission often tied to submissions
            'default_columns': [ "id", "form_submission_id","form_submission.form.title", "question", "question_type", "answer", "created_at" ],
            'analysis_hints': {
                'date_columns': ['created_at', 'updated_at'],
                'categorical_columns': ['question', 'question_type', 'form_submission.form.title'],
            },
            'default_sort': [{"field": "created_at", "direction": "desc"}],
            'stats_generators': ['generate_answers_submitted_stats', '_generate_generic_stats'],
            'chart_generators': ['_generate_generic_charts'], # Primarily generic charts
            'insight_generators': ['_generate_generic_insights'],
            'format_generators': {}
        },
        "attachments": {
            'model': Attachment, 'view_permission_entity': EntityType.ATTACHMENTS,
            'default_columns': [ "id", "form_submission_id", "form_submission.form.title", "file_name", "file_type", "is_signature", "signature_author", "created_at" ],
            'analysis_hints': {
                'date_columns': ['created_at', 'updated_at'],
                'categorical_columns': ['file_type', 'is_signature', 'signature_author', 'form_submission.form.title', 'file_name'],
            },
            'default_sort': [{"field": "created_at", "direction": "desc"}],
            'stats_generators': ['generate_attachment_stats', '_generate_generic_stats'],
            'chart_generators': ['_generate_generic_charts'], # Primarily generic charts
            'insight_generators': ['_generate_generic_insights'],
            'format_generators': {}
        },
        "permissions": { # Typically admin/dev focused, generic analysis is fine
            'model': Permission, 'view_permission_entity': EntityType.ROLES, # Permissions often viewed in context of roles
            'default_columns': ["id", "name", "action", "entity"],
            'analysis_hints': {'categorical_columns': ['name', 'action', 'entity']},
            'default_sort': [{"field": "name", "direction": "asc"}],
            'stats_generators': ['_generate_generic_stats'],
            'chart_generators': ['_generate_generic_charts'],
            'insight_generators': ['_generate_generic_insights'],
            'format_generators': {}
        },
        "role_permissions": { # Admin/dev focused
            'model': RolePermission, 'view_permission_entity': EntityType.ROLES,
            'default_columns': ["id", "role_id", "permission_id", "role.name", "permission.name", "permission.action", "permission.entity"],
            'analysis_hints': {'categorical_columns': ['role.name', 'permission.name', 'permission.action', 'permission.entity']},
            'default_sort': [{"field": "role_id", "direction": "asc"}, {"field": "permission_id", "direction": "asc"}],
            'stats_generators': ['_generate_generic_stats'],
            'chart_generators': ['_generate_generic_charts'],
            'insight_generators': ['_generate_generic_insights'],
            'format_generators': {}
        },
        "question_types": { # Reference data, generic analysis is fine
            'model': QuestionType, 'view_permission_entity': EntityType.QUESTION_TYPES,
            'default_columns': ["id", "type"],
            'analysis_hints': {'categorical_columns': ['type']},
            'default_sort': [{"field": "type", "direction": "asc"}],
            'stats_generators': ['_generate_generic_stats'],
            'chart_generators': ['_generate_generic_charts'],
            'insight_generators': ['_generate_generic_insights'],
            'format_generators': {}
        },
        "questions": { # Reference data, generic analysis suitable
            'model': Question, 'view_permission_entity': EntityType.QUESTIONS,
            'default_columns': ["id", "text", "question_type.type", "is_signature", "question_type_id"],
            'analysis_hints': {'categorical_columns': ['question_type.type', 'is_signature', 'text_shortened']}, # Assuming text_shortened could be a property or transformation
            'default_sort': [{"field": "text", "direction": "asc"}],
            'stats_generators': ['_generate_generic_stats'],
            'chart_generators': ['_generate_generic_charts'],
            'insight_generators': ['_generate_generic_insights'],
            'format_generators': {}
        },
        "answers": { # Predefined answers, generic analysis suitable
            'model': Answer, 'view_permission_entity': EntityType.ANSWERS,
            'default_columns': ["id", "value", "remarks"],
            'analysis_hints': {'categorical_columns': ['value_shortened']}, # Assuming value_shortened for long answers
            'default_sort': [{"field": "value", "direction": "asc"}],
            'stats_generators': ['_generate_generic_stats'],
            'chart_generators': ['_generate_generic_charts'],
            'insight_generators': ['_generate_generic_insights'],
            'format_generators': {}
        },
        "form_questions": { # Junction table, generic analysis often helpful
            'model': FormQuestion, 'view_permission_entity': EntityType.FORMS,
            'default_columns': ["id", "form_id", "question_id", "order_number", "is_required", "form.title", "question.text"],
            'analysis_hints': {'categorical_columns': ['form.title', 'question.text', 'is_required']},
            'default_sort': [{"field": "form_id", "direction": "asc"}, {"field": "order_number", "direction":"asc"}],
            'stats_generators': ['_generate_generic_stats'],
            'chart_generators': ['_generate_generic_charts'],
            'insight_generators': ['_generate_generic_insights'],
            'format_generators': {}
        },
        "form_answers": { # Junction table for predefined answers, generic analysis
            'model': FormAnswer, 'view_permission_entity': EntityType.FORMS,
            'default_columns': ["id", "form_question_id", "answer_id", "form_question.question.text", "answer.value"],
            'analysis_hints': {'categorical_columns': ['form_question.question.text', 'answer.value']},
            'default_sort': [{"field": "form_question_id", "direction": "asc"}],
            'stats_generators': ['_generate_generic_stats'],
            'chart_generators': ['_generate_generic_charts'],
            'insight_generators': ['_generate_generic_insights'],
            'format_generators': {}
        },
    }

    # --- Core Helper Methods ---
    @staticmethod
    def _get_attribute_recursive(obj: Any, attr_string: str) -> Any:
        """Recursively retrieves nested attributes, formatting specific types."""
        value = obj
        try:
            for attr in attr_string.split('.'):
                if value is None: return None
                if '[' in attr and attr.endswith(']'):
                    attr_name, index_str = attr.split('[', 1)
                    index = int(index_str[:-1])
                    list_attr = getattr(value, attr_name, None)
                    value = list_attr[index] if isinstance(list_attr, list) and len(list_attr) > index else None
                else:
                    value = getattr(value, attr, None)
            if isinstance(value, datetime): return value.isoformat(sep=' ', timespec='seconds')
            if isinstance(value, date): return value.isoformat()
            if isinstance(value, bool): return "Yes" if value else "No" # Keep boolean formatting consistent
            return value
        except (AttributeError, ValueError, IndexError, TypeError):
            return None

    # @staticmethod
    # def _get_sqlalchemy_attribute(base_model: Type, field_path: str, query: Query, joined_models: set, is_filter: bool) -> Tuple[Optional[Type], Optional[Any]]:
    #     """Resolves model attributes and handles joins for filters/sort."""
    #     current_model = base_model
    #     current_alias = None
    #     parts = field_path.split('.')
    #     for i, part in enumerate(parts):
    #         mapper = sqla_inspect(current_model)
    #         is_last_part = (i == len(parts) - 1)
    #         if is_last_part:
    #             if part in mapper.columns or part in mapper.synonyms or hasattr(current_model, part):
    #                 attr_source = current_alias if current_alias else current_model
    #                 if hasattr(attr_source, part):
    #                     potential_attr = getattr(attr_source, part)
    #                     return current_model, potential_attr
    #                 else:
    #                     logger.warning(f"Attribute '{part}' found in mapper but not on model/alias '{getattr(attr_source, '__name__', str(attr_source))}'")
    #                     return current_model, None
    #             else:
    #                 logger.warning(f"Attribute '{part}' not found on model {current_model.__name__}")
    #                 return current_model, None
    #         else:
    #             if part in mapper.relationships:
    #                 relationship = mapper.relationships[part]
    #                 related_model = relationship.mapper.class_
    #                 relationship_key = relationship.key
    #                 related_alias = aliased(related_model)
    #                 if current_alias:
    #                     query = query.join(related_alias, getattr(current_alias, relationship_key))
    #                 else:
    #                     query = query.join(related_alias, getattr(current_model, relationship_key))
    #                 current_model = related_model
    #                 current_alias = related_alias
    #             else:
    #                 logger.warning(f"Relationship '{part}' not found on model {current_model.__name__}")
    #                 return current_model, None
    #     return current_model, None

    @staticmethod
    def _apply_between_filter(model_attr: Any, field_name: str, value: List) -> Optional[Any]:
        """Applies a 'between' filter, attempting date/datetime conversion."""
        if not isinstance(value, list) or len(value) != 2: return None
        try:
            start, end = value[0], value[1]
            col_type = getattr(model_attr, 'type', None)
            is_date_type = isinstance(col_type, (db.Date, db.DateTime)) if col_type and hasattr(db, 'Date') else False
            is_likely_date = any(sub in field_name.lower() for sub in ["at", "date"])
            if is_date_type or is_likely_date:
                start_dt = pd.to_datetime(start, errors='coerce'); end_dt = pd.to_datetime(end, errors='coerce')
                if pd.notna(start_dt) and pd.notna(end_dt):
                    if isinstance(col_type, db.Date) and not isinstance(col_type, db.DateTime):
                         start_dt = start_dt.date(); end_dt = end_dt.date()
                    return model_attr.between(start_dt, end_dt)
                else: return None
            else: return model_attr.between(start, end)
        except Exception as e: logger.warning(f"Error applying 'between' filter on {field_name}: {e}"); return None

    @staticmethod
    def _apply_filters_and_sort(query: Query, model_cls: type, filters: List[Dict], sort_by: List[Dict]) -> Query:
        """
        Applies filtering and sorting to a SQLAlchemy query object, ensuring
        necessary joins are added correctly using aliases.
        """
        # Keep track of aliases used for joins to avoid redundant joins
        # Key: relationship path string (e.g., "form_submission"), Value: alias object
        join_aliases: Dict[str, Any] = {}

        # --- Apply Filters ---
        if filters:
            logger.debug(f"Applying filters for {model_cls.__name__}: {filters}")
            for f_idx, f in enumerate(filters):
                field, op, value = f.get("field"), f.get("operator", "eq").lower(), f.get("value")
                if not field or (value is None and op not in ["isnull", "isnotnull"]):
                    logger.debug(f"Skipping filter #{f_idx} (missing field/value): {f}")
                    continue

                logger.debug(f"Processing filter #{f_idx}: {field} {op} {value}")
                try:
                    # Resolve attribute and add necessary joins FOR FILTERING
                    # Pass the current state of the query and aliases
                    query, model_attr, join_aliases = ReportService._resolve_attribute_and_joins(
                        model_cls, field, query, join_aliases
                    )
                    if model_attr is None:
                        logger.warning(f"Could not resolve filter field: {field}")
                        continue

                    # Apply the filter condition (existing logic)
                    op_map = {
                        "eq": lambda a, v: a == v,
                        "neq": lambda a, v: a != v,
                        "like": lambda a, v: a.ilike(f"%{str(v)}%"),
                        "notlike": lambda a, v: ~a.ilike(f"%{str(v)}%"),
                        "in": lambda a, v: a.in_(v) if isinstance(v, list) else None,
                        "notin": lambda a, v: ~a.in_(v) if isinstance(v, list) else None,
                        "gt": lambda a, v: a > v,
                        "lt": lambda a, v: a < v,
                        "gte": lambda a, v: a >= v,
                        "lte": lambda a, v: a <= v,
                        "between": lambda a, v: ReportService._apply_between_filter(a, field, v),
                        "isnull": lambda a, v: a == None,
                        "isnotnull": lambda a, v: a != None,
                    }
                    filter_func = op_map.get(op)
                    if filter_func:
                        condition = None
                        if op in ["isnull", "isnotnull"]:
                             condition = filter_func(model_attr, None)
                        else:
                            # Handle boolean string conversion
                            if isinstance(getattr(model_attr, 'type', None), Boolean) and isinstance(value, str):
                                bool_map = {'true': True, 'false': False, 'yes': True, 'no': False, '1': True, '0': False}
                                value_lower = value.lower()
                                if value_lower in bool_map:
                                    value = bool_map[value_lower]
                                else:
                                    logger.warning(f"Could not interpret '{value}' as boolean for field '{field}'.")
                                    continue # Skip filter if boolean conversion fails
                            condition = filter_func(model_attr, value)

                        if condition is not None:
                            query = query.filter(condition)
                            logger.debug(f"Applied filter #{f_idx}: {field} {op} {value}")
                        else:
                            logger.warning(f"Invalid filter value/op for filter #{f_idx} ('{op}' on field '{field}'): {value}")
                    else:
                        logger.warning(f"Unsupported filter operator '{op}' for filter #{f_idx} on field '{field}'")

                except Exception as e:
                    logger.warning(f"Could not apply filter #{f_idx} ({f}): {e}", exc_info=True)

        # --- Apply Sorting ---
        if sort_by:
            logger.debug(f"Applying sorting for {model_cls.__name__}: {sort_by}")
            for s_idx, s in enumerate(sort_by):
                field, direction = s.get("field"), s.get("direction", "asc").lower()
                if not field:
                    logger.debug(f"Skipping sort #{s_idx} (missing field)")
                    continue
                if direction not in ["asc", "desc"]:
                    logger.warning(f"Invalid sort direction '{direction}' for sort #{s_idx} on field '{field}'. Defaulting to 'asc'.")
                    direction = "asc"

                logger.debug(f"Processing sort #{s_idx}: {field} {direction.upper()}")
                try:
                    # Resolve attribute and add necessary joins FOR SORTING
                    # Pass the potentially modified query from filtering stage and current aliases
                    query, sort_attr, join_aliases = ReportService._resolve_attribute_and_joins(
                        model_cls, field, query, join_aliases
                    )

                    if sort_attr is not None:
                        order_func = desc if direction == "desc" else asc
                        query = query.order_by(order_func(sort_attr))
                        logger.debug(f"Applied sort #{s_idx}: {field} {direction.upper()}")
                    else:
                        logger.warning(f"Could not resolve sort field for sort #{s_idx}: {field}")
                except Exception as e:
                    logger.error(f"Error applying sort #{s_idx} for field '{field}': {e}", exc_info=True)

        return query # Return the query with filters and sorts applied

    @staticmethod
    def _apply_rbac_filters(query: Query, model_cls: type, user: User) -> Query:
        # (Keep existing implementation)
        if user.role and user.role.is_super_user: return query
        env_id = user.environment_id; user_role_name = user.role.name if user.role else None
        if model_cls == User: return query.filter(User.environment_id == env_id)
        elif model_cls == FormSubmission:
            FormCreatorUser = aliased(User)
            if user_role_name in [RoleType.SITE_MANAGER, RoleType.SUPERVISOR]: return query.join(Form, Form.id == FormSubmission.form_id).join(FormCreatorUser, FormCreatorUser.id == Form.user_id).filter(FormCreatorUser.environment_id == env_id)
            else: return query.filter(FormSubmission.submitted_by == user.username)
        elif model_cls == Form: FormCreatorUser = aliased(User); return query.join(FormCreatorUser, FormCreatorUser.id == Form.user_id).filter(or_(Form.is_public == True, FormCreatorUser.environment_id == env_id))
        elif model_cls == Environment: return query.filter(Environment.id == env_id)
        elif model_cls == Role: return query.filter(Role.is_super_user == False)
        elif model_cls in [AnswerSubmitted, Attachment]:
             link_relationship = getattr(model_cls, 'form_submission', None)
             if link_relationship is None: logger.error(f"RBAC Error: {model_cls.__name__}"); return query.filter(False)
             FormCreatorUser = aliased(User); return query.join(link_relationship).join(Form, Form.id == FormSubmission.form_id).join(FormCreatorUser, FormCreatorUser.id == Form.user_id).filter(FormCreatorUser.environment_id == env_id)
        elif model_cls == RolePermission: return query.join(Role, Role.id == RolePermission.role_id).filter(Role.is_super_user == False)
        elif model_cls == FormQuestion: FormCreatorUser = aliased(User); return query.join(Form, Form.id == FormQuestion.form_id).join(FormCreatorUser, FormCreatorUser.id == Form.user_id).filter(or_(Form.is_public == True, FormCreatorUser.environment_id == env_id))
        elif model_cls == FormAnswer: FormCreatorUser = aliased(User); return query.join(FormQuestion, FormQuestion.id == FormAnswer.form_question_id).join(Form, Form.id == FormQuestion.form_id).join(FormCreatorUser, FormCreatorUser.id == Form.user_id).filter(or_(Form.is_public == True, FormCreatorUser.environment_id == env_id))
        if hasattr(model_cls, 'environment_id'): logger.warning(f"Applying default env RBAC for {model_cls.__name__}"); return query.filter(model_cls.environment_id == env_id)
        elif hasattr(model_cls, 'user_id'): logger.warning(f"Applying default user RBAC for {model_cls.__name__}"); return query.filter(model_cls.user_id == user.id)
        logger.warning(f"No specific RBAC rule for {model_cls.__name__} for role {user_role_name}. Allowing.")
        return query@staticmethod
    def _get_load_options(model_cls: Type, requested_columns: List[str]) -> list:
        # (Keep existing implementation)
        options = []; relationship_paths = set()
        for col in requested_columns:
            if col.startswith(ANSWERS_PREFIX) and model_cls == FormSubmission: relationship_paths.add('answers_submitted'); continue
            parts = col.split('.');
            if len(parts) > 1:
                for i in range(1, len(parts)): relationship_paths.add('.'.join(parts[:i]))
        processed_base_paths = set(); sorted_paths = sorted(list(relationship_paths), key=lambda x: (x.count('.'), x), reverse=True)
        for path_str in sorted_paths:
            parts = path_str.split('.'); base_path = parts[0]
            if base_path in processed_base_paths and len(parts) == 1: continue
            current_load_option = None; current_model = model_cls
            try:
                for i, part in enumerate(parts):
                    mapper = sqla_inspect(current_model)
                    if part not in mapper.relationships: logger.warning(f"Invalid relationship '{part}' in path '{path_str}' for {current_model.__name__}."); current_load_option = None; break
                    relationship = mapper.relationships[part]; load_func = selectinload if relationship.uselist else joinedload
                    if i == 0: current_load_option = load_func(getattr(current_model, part))
                    elif current_load_option: current_load_option = current_load_option.selectinload(getattr(current_model, part)) if relationship.uselist else current_load_option.joinedload(getattr(current_model, part))
                    current_model = relationship.mapper.class_
                if current_load_option is not None: options.append(current_load_option); processed_base_paths.add(base_path)
            except AttributeError as ae: logger.error(f"AttributeError building load option for '{path_str}': {ae}")
            except Exception as e: logger.error(f"Error building load option for '{path_str}': {e}", exc_info=True)
        return options

    @staticmethod
    def _fetch_data(model_cls: type, filters: List[Dict], sort_by: List[Dict], user: User, requested_columns: List[str]) -> List[Any]:
        # (Keep existing implementation)
        try:
            query = db.session.query(model_cls)
            if hasattr(model_cls, 'is_deleted'): query = query.filter(model_cls.is_deleted == False)
            query = ReportService._apply_rbac_filters(query, model_cls, user)
            query = ReportService._apply_filters_and_sort(query, model_cls, filters, sort_by)
            load_options = ReportService._get_load_options(model_cls, requested_columns)
            if load_options: query = query.options(*load_options)
            results = query.all()
            logger.info(f"Fetched {len(results)} records for {model_cls.__name__}.")
            return results
        except Exception as e: logger.error(f"Error fetching data for {model_cls.__name__}: {e}", exc_info=True); raise

    @staticmethod
    def _flatten_data(objects: List[Any], columns: List[str], report_type: str) -> DataList:
        # (Keep existing implementation)
        if not objects: return []
        flat_data = []; question_info_map = {}
        if report_type == 'form_submissions':
             all_questions_text = set()
             for obj in objects:
                 if hasattr(obj, 'answers_submitted'):
                     for ans_sub in obj.answers_submitted:
                         if ans_sub and not ans_sub.is_deleted: all_questions_text.add(ans_sub.question)
             if all_questions_text:
                  questions_from_db = db.session.query(Question.text, QuestionType.type).join(QuestionType, Question.question_type_id == QuestionType.id).filter(Question.text.in_(list(all_questions_text))).all()
                  question_info_map = {q_text: q_type for q_text, q_type in questions_from_db}
        for obj in objects:
            row_dict = {}
            dynamic_columns_requested = False
            for col in columns:
                 if col.startswith(ANSWERS_PREFIX) and report_type == 'form_submissions': dynamic_columns_requested = True; continue
                 row_dict[col] = ReportService._get_attribute_recursive(obj, col)
            if dynamic_columns_requested and report_type == 'form_submissions' and hasattr(obj, 'answers_submitted'):
                submission_answers = {ans.question: ans.answer for ans in obj.answers_submitted if ans and not ans.is_deleted}
                for col in columns:
                    if col.startswith(ANSWERS_PREFIX):
                        try: question_text = col.split(ANSWERS_PREFIX, 1)[1]; answer_value = submission_answers.get(question_text); row_dict[col] = answer_value
                        except IndexError: logger.warning(f"Invalid dynamic answer column format: {col}"); row_dict[col] = None
            flat_data.append(row_dict)
        return flat_data

    # --- Data Analysis ---
    @staticmethod
    def _analyze_data(data: DataList, params: ParamsDict, report_type: str) -> AnalysisDict:
        # (Keep existing implementation)
        analysis: AnalysisDict = {"summary_stats": {}, "charts": {}, "insights": {}}
        if not data: analysis['insights']['status'] = "No data available for analysis."; return analysis
        try:
            df = pd.DataFrame(data)
            config = params.get('_internal_config', {})
            analysis['summary_stats']['record_count'] = len(df)
            analysis['_internal_params'] = params
            analysis['_internal_config'] = config
            hints = config.get('analysis_hints', {})
            for col in hints.get('date_columns', []):
                if col in df.columns:
                    try: df[col] = pd.to_datetime(df[col], errors='coerce')
                    except Exception as date_err: logger.warning(f"Could not parse date column '{col}' for {report_type}: {date_err}")
            if report_type == 'form_submissions' and ANSWERS_PREFIX in hints.get('dynamic_answer_prefix', ''):
                 question_info_map = params.get('_internal_question_info', {})
                 for col in df.columns:
                     if col.startswith(ANSWERS_PREFIX):
                         question_text = col.split(ANSWERS_PREFIX, 1)[1]; q_type = question_info_map.get(question_text)
                         if q_type in ['date', 'datetime']:
                             try: df[col] = pd.to_datetime(df[col], errors='coerce')
                             except Exception as date_err: logger.warning(f"Could not parse dynamic date column '{col}': {date_err}")
                         elif q_type in ['integer', 'number']:
                             try: df[col] = pd.to_numeric(df[col], errors='coerce')
                             except Exception as num_err: logger.warning(f"Could not parse dynamic numeric column '{col}': {num_err}")
            analysis['_internal_df'] = df
            for gen_type in ['stats', 'chart', 'insight']:
                 key = f"{gen_type}_generators"; output_key = f"summary_stats" if gen_type == 'stats' else f"{gen_type}s"
                 for func_name in config.get(key, []):
                     if hasattr(ReportService, func_name):
                         try:
                             generator_func = getattr(ReportService, func_name)
                             if gen_type == 'stats': args = (df, params)
                             elif gen_type == 'chart': args = (df, analysis) # Pass analysis for context
                             elif gen_type == 'insight': args = (df, analysis, params)
                             else: args = (df, params)
                             result = generator_func(*args)
                             if result and isinstance(result, dict): analysis[output_key].update(result)
                         except Exception as gen_err: logger.error(f"Error executing {gen_type} function '{func_name}': {gen_err}", exc_info=True); analysis[output_key][f'{func_name}_error'] = f"Failed: {gen_err}"
        except Exception as e: logger.error(f"Error during data analysis for {report_type}: {e}", exc_info=True); analysis['error'] = f"Analysis failed: {e}"; analysis['insights']['status'] = "Error during data analysis."
        analysis.pop('_internal_df', None); analysis.pop('_internal_params', None); analysis.pop('_internal_config', None)
        return analysis

    @staticmethod
    def _safe_value_counts(df: pd.DataFrame, column: str, top_n: Optional[int] = 10) -> Dict:
        # (Keep existing implementation)
        if column not in df.columns or df[column].isnull().all(): return {}
        try:
            if df[column].apply(lambda x: isinstance(x, list)).any(): counts = df[column].apply(lambda x: tuple(x) if isinstance(x, list) else x).value_counts()
            else: counts = df[column].value_counts()
            if top_n: counts = counts.nlargest(top_n)
            return {str(k): int(v) if isinstance(v, np.integer) else float(v) if isinstance(v, np.floating) else v for k, v in counts.to_dict().items()}
        except Exception as e: logger.warning(f"Could not calculate value counts for column '{column}': {e}"); return {}

    # --- Specific Stats Generators ---
    @staticmethod
    def generate_submission_stats(df: pd.DataFrame, params: ParamsDict) -> Dict:
        # (Keep existing implementation)
        stats = {}; question_info = params.get('_internal_question_info', {})
        for col in df.columns:
            if col.startswith(ANSWERS_PREFIX):
                 question_text = col.split(ANSWERS_PREFIX, 1)[1]; q_type = question_info.get(question_text)
                 if q_type in ['multiple_choices', 'dropdown', 'user', 'checkbox']:
                     col_stats = ReportService._safe_value_counts(df, col, top_n=10);
                     if col_stats: stats[f'counts_{question_text.lower().replace(" ", "_").replace("?", "")[:30]}'] = col_stats
                 elif q_type in ['date', 'datetime'] and pd.api.types.is_datetime64_any_dtype(df[col]):
                      valid_dates = df[col].dropna();
                      if not valid_dates.empty: stats[f'range_{question_text.lower().replace(" ", "_").replace("?", "")[:30]}'] = {'first': valid_dates.min().isoformat(), 'last': valid_dates.max().isoformat()}
        submitter_stats = ReportService._safe_value_counts(df, 'submitted_by', top_n=5);
        if submitter_stats: stats['submissions_per_user_top5'] = submitter_stats
        form_stats = ReportService._safe_value_counts(df, 'form.title', top_n=5);
        if form_stats: stats['submissions_per_form_top5'] = form_stats
        if 'submitted_at' in df.columns and pd.api.types.is_datetime64_any_dtype(df['submitted_at']):
             valid_dates = df['submitted_at'].dropna()
             if not valid_dates.empty:
                 stats['overall_submission_range'] = {'first': valid_dates.min().isoformat(), 'last': valid_dates.max().isoformat()}
                 date_range_days = (valid_dates.max() - valid_dates.min()).days
                 if date_range_days is not None and date_range_days >= 0: stats['average_daily_submissions'] = round(len(df) / max(date_range_days, 1), 1)
                 else: stats['average_daily_submissions'] = len(df)
        return stats

    @staticmethod
    def generate_user_stats(df: pd.DataFrame, params: ParamsDict) -> Dict:
        # (Keep existing implementation)
        stats = {}; stats['users_per_role'] = ReportService._safe_value_counts(df, 'role.name'); stats['users_per_environment'] = ReportService._safe_value_counts(df, 'environment.name')
        if 'created_at' in df.columns and pd.api.types.is_datetime64_any_dtype(df['created_at']):
            valid_dates = df['created_at'].dropna();
            if not valid_dates.empty: stats['user_creation_range'] = {'first': valid_dates.min().isoformat(), 'last': valid_dates.max().isoformat()}
        return stats

    @staticmethod
    def generate_form_stats(df: pd.DataFrame, params: ParamsDict) -> Dict:
       # (Keep existing implementation)
       stats = {}; stats['forms_per_creator_top5'] = ReportService._safe_value_counts(df, 'creator.username', top_n=5); stats['forms_per_environment'] = ReportService._safe_value_counts(df, 'creator.environment.name')
       if 'is_public' in df.columns: stats['public_vs_private_forms'] = ReportService._safe_value_counts(df, 'is_public')
       return stats

    @staticmethod
    def generate_environment_stats(df: pd.DataFrame, params: ParamsDict) -> Dict:
       # (Keep existing implementation)
       return {'total_environments_reported': len(df)}

    @staticmethod
    def generate_role_stats(df: pd.DataFrame, params: ParamsDict) -> Dict:
       # (Keep existing implementation)
       stats = {}
       if 'is_super_user' in df.columns: stats['roles_by_superuser_status'] = ReportService._safe_value_counts(df, 'is_super_user')
       return stats

    @staticmethod
    def generate_answers_submitted_stats(df: pd.DataFrame, params: ParamsDict) -> Dict:
        # (Keep existing implementation)
        stats = {}; stats['answers_per_question_type'] = ReportService._safe_value_counts(df, 'question_type', top_n=10); stats['answers_per_question_text_top10'] = ReportService._safe_value_counts(df, 'question', top_n=10)
        return stats

    @staticmethod
    def generate_attachment_stats(df: pd.DataFrame, params: ParamsDict) -> Dict:
        # (Keep existing implementation)
        stats = {}; stats['attachments_by_type'] = ReportService._safe_value_counts(df, 'file_type')
        if 'is_signature' in df.columns: stats['attachments_by_signature_status'] = ReportService._safe_value_counts(df, 'is_signature')
        stats['attachments_per_author_top5'] = ReportService._safe_value_counts(df, 'signature_author', top_n=5)
        return stats

    # --- Generic Stats Generator ---
    @staticmethod
    def _generate_generic_stats(df: pd.DataFrame, params: ParamsDict) -> Dict[str, Any]:
        stats: Dict[str, Any] = {}
        report_type = params.get("report_type", "unknown_entity")
        logger.info(f"Generating generic stats for {report_type}")

        if df.empty:
            stats['record_count'] = 0
            return stats

        stats['record_count'] = len(df)
        
        # --- Generic Categorical Stats ---
        categorical_stats_count = 0
        MAX_GENERIC_CATEGORICAL_STATS = 5 # Limit the number of distinct categorical columns we generate stats for

        # Iterate through all columns present in the DataFrame `df`
        # (which was built from user's requested columns or defaults)
        for col_name in df.columns:
            if categorical_stats_count >= MAX_GENERIC_CATEGORICAL_STATS:
                break

            # Ensure col_name is treated as a string for operations, though df.columns usually are.
            current_col_name_str = str(col_name)

            is_suitable_categorical = False
            # Use dropna=True so that uniqueness is not affected by presence of NaNs for decision-making
            num_unique = df[current_col_name_str].nunique(dropna=True) 
            
            # Define a range for what's considered a "countable" or "categorical" column for stats
            # We allow more unique values for stats (top_n will limit) than for direct charting.
            is_in_stat_range = 1 < num_unique <= (MAX_UNIQUE_GENERIC_CHART * 2) # MAX_UNIQUE_GENERIC_CHART is 15, so 2-30 unique values here.

            col_dtype = df[str(col_name)].dtype # Ensure col_name is string for df access
            current_col_name_str = str(col_name)

            is_suitable_categorical = False
            num_unique = df[current_col_name_str].nunique(dropna=True)

            if pd.api.types.is_bool_dtype(col_dtype):
                is_suitable_categorical = True
            elif pd.api.types.is_integer_dtype(col_dtype):
                # For integers (like IDs), keep the stricter unique check
                if 1 < num_unique <= (MAX_UNIQUE_GENERIC_CHART * 2): # e.g., 2-30
                    is_suitable_categorical = True
            elif pd.api.types.is_string_dtype(col_dtype) or col_dtype == 'object':
                # For strings/objects (like username), be more lenient on total unique count
                # as _safe_value_counts will take top_n. We just need > 1 unique value.
                if 1 < num_unique: # If there's more than one unique string, consider it
                    is_suitable_categorical = True
            
            if is_suitable_categorical:
                col_stats_values = ReportService._safe_value_counts(df, current_col_name_str, top_n=10)
                if col_stats_values:
                    safe_col_key = ''.join(c if c.isalnum() else '_' for c in current_col_name_str)[:30]
                    stats[f'counts_{safe_col_key}'] = col_stats_values
                    categorical_stats_count += 1
                    logger.debug(f"Generated generic count stat for column '{current_col_name_str}' in {report_type}")

        # --- Generic Date Range Stats ---
        # This logic prioritizes hinted date columns but also scans the DataFrame.
        hinted_date_cols = params.get('_internal_config', {}).get('analysis_hints', {}).get('date_columns', [])
        date_candidates = set(hinted_date_cols)
        
        for col_name in df.columns: # Also check all df columns for datetime types
            if pd.api.types.is_datetime64_any_dtype(df[str(col_name)].dtype): # Check dtype directly
                date_candidates.add(str(col_name))
        
        # Generate stat for the first valid date column found
        for col_name_str in list(date_candidates): # Iterate a copy
            if col_name_str in df.columns and pd.api.types.is_datetime64_any_dtype(df[col_name_str].dtype):
                valid_dates = df[col_name_str].dropna()
                if not valid_dates.empty:
                    safe_col_key = ''.join(c if c.isalnum() else '_' for c in col_name_str)[:30]
                    stats[f'range_{safe_col_key}'] = {
                        'first': valid_dates.min().isoformat(),
                        'last': valid_dates.max().isoformat()
                    }
                    logger.debug(f"Generated generic date range stat for column '{col_name_str}' in {report_type}")
                    break # Only take the first date column found for a generic range stat

        logger.debug(f"Generated generic stats for {report_type}: {list(stats.keys())}")
        return stats

    # --- Insight Generators (Specific - Unchanged) ---
    @staticmethod
    def generate_submission_insights(df: pd.DataFrame, analysis: AnalysisDict, params: ParamsDict) -> Dict[str, str]:
        # (Keep existing implementation)
        insights = {}; stats = analysis.get('summary_stats', {}); record_count = stats.get('record_count', 0)
        if record_count == 0: return {"status": "No submission data available for analysis."}
        insights["volume"] = f"Analyzed {record_count} total submissions."
        if stats.get('overall_submission_range'): first = stats['overall_submission_range']['first'].split(' ')[0]; last = stats['overall_submission_range']['last'].split(' ')[0]; insights["date_range"] = f"Data spans from {first} to {last}."
        if stats.get('average_daily_submissions'): insights["activity_rate"] = f"Average daily submission rate: {stats['average_daily_submissions']:.1f}."
        if stats.get('submissions_per_user_top5'): top_user = next(iter(stats['submissions_per_user_top5']), 'N/A'); insights["top_user"] = f"The most active user was '{top_user}'."
        dept_stats_key = next((k for k in stats if k.startswith('counts_') and 'department' in k), None)
        if dept_stats_key and stats[dept_stats_key]: top_dept = next(iter(stats[dept_stats_key]), 'N/A'); insights["top_department"] = f"'{top_dept}' submitted the most forms."
        return insights

    @staticmethod
    def generate_user_insights(df: pd.DataFrame, analysis: AnalysisDict, params: ParamsDict) -> Dict[str, str]:
        # (Keep existing implementation)
        insights = {}; stats = analysis.get('summary_stats', {}); record_count = stats.get('record_count', 0)
        if record_count == 0: return {"status": "No user data available for analysis."}
        insights["user_count"] = f"Analyzed {record_count} user records."
        if stats.get('users_per_role'): insights["role_distribution"] = f"Users are distributed across {len(stats['users_per_role'])} roles."
        if stats.get('users_per_environment'): insights["env_distribution"] = f"Users belong to {len(stats['users_per_environment'])} environments."
        return insights

    @staticmethod
    def generate_form_insights(df: pd.DataFrame, analysis: AnalysisDict, params: ParamsDict) -> Dict[str, str]:
        # (Keep existing implementation)
        insights = {}; stats = analysis.get('summary_stats', {}); record_count = stats.get('record_count', 0)
        if record_count == 0: return {"status": "No form data available for analysis."}
        insights["form_count"] = f"Analyzed {record_count} form records."
        if stats.get('public_vs_private_forms'): public = stats['public_vs_private_forms'].get('Yes', 0); private = stats['public_vs_private_forms'].get('No', 0); insights["public_status"] = f"{public} public and {private} private forms found."
        if stats.get('forms_per_creator_top5'): insights["creator_activity"] = f"Top creators identified based on form count."
        return insights

    # --- Generic Insights Generator ---
    @staticmethod
    def _generate_generic_insights(df: pd.DataFrame, analysis: AnalysisDict, params: ParamsDict) -> Dict[str, str]:
        insights = {} # This function should return what IT generates
        stats = analysis.get('summary_stats', {})
        report_type = params.get("report_type", "this entity")
        record_count = stats.get('record_count', 0)

        if record_count == 0:
            # If 'status' isn't already set by a more specific insight generator
            if 'status' not in analysis.get('insights', {}):
                 insights["status"] = f"No data available for {report_type}."
            return insights # Return only the status if no data

        # Avoid duplicating record count if a specific insight generator already added it
        # This relies on convention that specific insights might use keys like 'user_count', 'form_count', 'volume'
        # and that they populate analysis['insights'] directly.
        # Note: _analyze_data updates analysis['insights'] by merging.
        existing_insights = analysis.get('insights', {})
        if not any(k in existing_insights for k in ['user_count', 'form_count', 'submission_count', 'record_summary', 'volume']):
            insights["record_summary"] = f"A total of {record_count} records were analyzed for {report_type}."
        
        # Date range insight (already improved)
        date_range_key = next((k for k in stats if k.startswith('range_')), None)
        if date_range_key and isinstance(stats[date_range_key], dict):
            first = stats[date_range_key].get('first', 'N/A')
            last = stats[date_range_key].get('last', 'N/A')
            date_col_name = date_range_key.replace('range_', '').replace('_', ' ')
            if first == last:
                insights["date_info"] = f"All analyzed records share the same timestamp for '{date_col_name}': {first.split('T')[0]}."
            else:
                insights["date_range"] = f"Records for '{date_col_name}' span from {first.split('T')[0]} to {last.split('T')[0]}."

        # Top category insight (already improved)
        count_keys = [k for k in stats if k.startswith('counts_')]
        if count_keys:
            # Sort keys to get a more consistent "primary" count key if multiple exist
            # Example: prefer 'counts_username' over 'counts_id' if both generated
            sorted_count_keys = sorted(count_keys, key=lambda k: (
                0 if 'name' in k or 'title' in k or 'type' in k else 1, # Prioritize common categorical names
                len(stats.get(k, {})), # Then by number of categories in the stat (fewer is often more focused)
                k # Alphabetical as tie-breaker
            ))
            primary_count_key = sorted_count_keys[0]
            
            if stats.get(primary_count_key): 
                counts_for_top_key = stats[primary_count_key]
                col_name = primary_count_key.replace('counts_', '').replace('_', ' ')
                if counts_for_top_key:
                    top_value = next(iter(counts_for_top_key), 'N/A')
                    top_count = counts_for_top_key.get(top_value, 0)
                    
                    is_dominant = True
                    if len(counts_for_top_key) > 1:
                        # Convert values to list for sorting if it's a dict
                        all_counts_in_stat = list(counts_for_top_key.values())
                        all_counts_in_stat.sort(reverse=True)
                        second_value_count = all_counts_in_stat[1]
                        if top_count <= second_value_count or top_count <= 1:
                            is_dominant = False
                    elif top_count <=1 and len(counts_for_top_key) == 1: # Only one category and its count is 1
                         is_dominant = False


                    if is_dominant:
                        insights["top_category_info"] = f"For the column '{col_name}', the most common value was '{top_value}', occurring {top_count} times."
                    elif top_count > 0 : 
                        insights["category_analysis_note"] = f"Distribution analysis was performed for '{col_name}'. The most frequent value found was '{top_value}' (occurrences: {top_count})."
        return insights

    # --- Charting Helpers & Generators (Specific) ---
    @staticmethod
    def _save_plot_to_bytes(figure) -> Optional[BytesIO]:
        # (Keep existing implementation)
        if not figure: return None
        try:
            img_buffer = BytesIO();
            try: figure.tight_layout(pad=1.1)
            except ValueError: logger.warning("Tight layout failed.")
            figure.savefig(img_buffer, format='png', dpi=150, bbox_inches='tight', facecolor=figure.get_facecolor())
            img_buffer.seek(0); plt.close(figure); return img_buffer
        except Exception as e: logger.error(f"Failed to save plot: {e}"); return None
        finally: plt.close('all')

    @staticmethod
    def _setup_chart(figsize=(10, 5), style='seaborn-v0_8-whitegrid') -> Tuple[plt.Figure, plt.Axes]:
        # (Keep existing implementation)
        plt.style.use(style); fig, ax = plt.subplots(figsize=figsize, facecolor='white'); ax.set_facecolor('white')
        plt.rcParams.update({'font.size': 10, 'axes.titlesize': 14, 'axes.labelsize': 12, 'xtick.labelsize': 10, 'ytick.labelsize': 10, 'legend.fontsize': 10})
        return fig, ax

    @staticmethod
    def _add_bar_labels(ax: plt.Axes, bars: Any, location: str = 'edge', color: str = 'black', fontsize=9, **kwargs):
        # (Keep existing implementation)
        try:
            patches = bars.patches if hasattr(bars, 'patches') else bars
            for bar in patches:
                x_offset = 0; y_offset = 0; height = bar.get_height(); width = bar.get_width()
                if height == 0 and width == 0: continue
                if height > 0: y_pos = height if location == 'edge' else height / 2; x_pos = bar.get_x() + width / 2; va = 'bottom' if location == 'edge' else 'center'; ha = 'center'; label_text = f'{int(height)}' if height == int(height) else f'{height:.1f}'; text_color = color if location != 'edge' else 'black'; y_offset = 3 if location == 'edge' else 0
                elif width > 0: x_pos = width if location == 'edge' else width / 2; y_pos = bar.get_y() + bar.get_height() / 2; ha = 'left' if location == 'edge' else 'center'; va = 'center'; label_text = f'{int(width)}' if width == int(width) else f'{width:.1f}'; text_color = color if location != 'edge' else 'black'; x_offset = 3 if location == 'edge' else 0
                else: continue
                ax.annotate(label_text, (x_pos, y_pos), xytext=(x_offset, y_offset), textcoords="offset points", ha=ha, va=va, color=text_color, fontsize=fontsize, fontweight='bold' if location == 'center' else 'normal', **kwargs)
        except Exception as e: logger.error(f"Error adding bar labels: {e}")

    @staticmethod
    def generate_submission_charts(df: pd.DataFrame, analysis: AnalysisDict) -> Dict[str, BytesIO]:
        # (Keep existing implementation)
        charts = {}; stats = analysis.get('summary_stats', {}); params = analysis.get('_internal_params', {})
        if df.empty: return charts
        config = analysis.get('_internal_config', {}); hints = config.get('analysis_hints', {}); date_col = next((c for c in hints.get('date_columns', []) if c == 'submitted_at'), None)
        if not date_col or date_col not in df.columns or not pd.api.types.is_datetime64_any_dtype(df[date_col]) or df[date_col].isnull().all(): logger.warning("Sub charts require valid 'submitted_at'"); return charts
        try:
            df_ts = df.set_index('submitted_at'); monthly_counts = df_ts.resample('ME').size()
            if not monthly_counts.empty and len(monthly_counts) > 1: fig, ax = ReportService._setup_chart(figsize=(10, 4)); monthly_counts.plot(kind='line', ax=ax, marker='o', color=sns.color_palette("viridis", 1)[0], linewidth=2); ax.set_title("Submissions Trend (Monthly)", fontsize=14, fontweight='bold'); ax.set_ylabel('# Submissions', fontsize=12); ax.set_xlabel(''); plt.xticks(rotation=30, ha='right'); plt.grid(axis='y', linestyle='--', alpha=0.6); charts['time_series_monthly'] = ReportService._save_plot_to_bytes(fig)
        except Exception as e: logger.error(f"Error generating monthly time series: {e}")
        try: charts['activity_heatmap'] = ReportService._create_activity_heatmap(df, date_col=date_col)
        except Exception as e: logger.error(f"Error generating activity heatmap: {e}")
        try:
            user_dist_stats = stats.get('submissions_per_user_top5');
            if user_dist_stats and isinstance(user_dist_stats, dict): user_counts = pd.Series(user_dist_stats);
            if not user_counts.empty: charts['user_distribution'] = ReportService._create_bar_chart(user_counts, 'Top 5 Users by Submissions', 'Username', '# Submissions', figsize=(8,4), palette='viridis')
        except Exception as e: logger.error(f"Error generating user dist chart: {e}")
        try:
            form_dist_stats = stats.get('submissions_per_form_top5');
            if form_dist_stats and isinstance(form_dist_stats, dict): form_counts = pd.Series(form_dist_stats);
            if not form_counts.empty: charts['form_distribution'] = ReportService._create_pie_chart( form_counts, 'Submissions by Form Type (Top 5)', figsize=(7, 5.5))
        except Exception as e: logger.error(f"Error generating form dist chart: {e}")
        try:
            for key, value_counts in stats.items():
                 if key.startswith('counts_') and isinstance(value_counts, dict):
                      question_text = key.replace('counts_', '').replace('_', ' ').title(); counts_series = pd.Series(value_counts)
                      if not counts_series.empty: chart_key = f'dist_{key.replace("counts_","")[:20]}'; charts[chart_key] = ReportService._create_bar_chart( counts_series, f'Distribution by: {question_text}', 'Answer', '# Submissions', figsize=(8,4), palette='viridis')
        except Exception as e: logger.error(f"Error generating dynamic charts from stats: {e}")
        return charts

    @staticmethod
    def generate_user_charts(df: pd.DataFrame, analysis: AnalysisDict) -> Dict[str, BytesIO]:
        charts = {}
        stats = analysis.get('summary_stats', {})
        if df.empty:
            return charts

        # Handle role distribution chart
        role_counts_data = stats.get('users_per_role')
        if role_counts_data:  # Check if the data exists in stats
            role_counts = pd.Series(role_counts_data)
            if not role_counts.empty:
                try:
                    charts['user_role_distribution'] = ReportService._create_bar_chart(
                        role_counts, 'User Count by Role', 'Role', '# Users', palette='Blues_d'
                    )
                except Exception as e:
                    logger.error(f"Error creating user_role_distribution chart: {e}", exc_info=True)
                    charts['user_role_distribution_error'] = f"Failed to generate: {e}"


        # Handle environment distribution chart
        env_counts_data = stats.get('users_per_environment')
        if env_counts_data:  # Check if the data exists in stats
            env_counts = pd.Series(env_counts_data)
            if not env_counts.empty:
                try:
                    charts['user_environment_distribution'] = ReportService._create_bar_chart(
                        env_counts, 'User Count by Environment', 'Environment', '# Users', palette='Greens_d'
                    )
                except Exception as e:
                    logger.error(f"Error creating user_environment_distribution chart: {e}", exc_info=True)
                    charts['user_environment_distribution_error'] = f"Failed to generate: {e}"
        return charts

    @staticmethod
    def generate_form_charts(df: pd.DataFrame, analysis: AnalysisDict) -> Dict[str, BytesIO]:
       # (Keep existing implementation)
       charts = {}; stats = analysis.get('summary_stats', {})
       if df.empty: return charts
       if stats.get('public_vs_private_forms'): status_counts = pd.Series(stats['public_vs_private_forms']);
       if not status_counts.empty: charts['form_public_private'] = ReportService._create_pie_chart( status_counts, 'Forms: Public vs. Private', figsize=(6,4))
       if stats.get('forms_per_creator_top5'): creator_counts = pd.Series(stats['forms_per_creator_top5']);
       if not creator_counts.empty: charts['forms_per_creator'] = ReportService._create_bar_chart( creator_counts, 'Top 10 Form Creators', 'Creator Username', '# Forms', palette='Oranges_d')
       return charts

    # --- Generic Chart Generator ---
    # In app/services/report_service.py

    @staticmethod
    def _generate_generic_charts(df: pd.DataFrame, analysis: AnalysisDict) -> Dict[str, BytesIO]:
        charts = {}
        stats = analysis.get('summary_stats', {})
        params = analysis.get('_internal_params', {}) # Assuming _analyze_data adds this
        report_type = params.get("report_type", "entity")
        logger.info(f"Generating generic charts for {report_type}")
        
        chart_count = 0 
        MAX_GENERIC_CHARTS = 3 

        for key, value_counts_dict in stats.items():
            if chart_count >= MAX_GENERIC_CHARTS:
                break
            
            if key.startswith('counts_') and isinstance(value_counts_dict, dict) and value_counts_dict:
                try:
                    counts_series = pd.Series(value_counts_dict)
                    col_name_from_key = key.replace('counts_', '').replace('_', ' ')
                    chart_title = f"Distribution by {col_name_from_key.title()}"
                    chart_key_name = f"generic_dist_{key.replace('counts_','')[:20].replace('.', '_')}"

                    if not counts_series.empty and len(counts_series.index) <= MAX_UNIQUE_GENERIC_CHART:
                        
                        # ---- START: Reverted Suppression Logic ----
                        # The following block was added to suppress flat charts.
                        # If you want to see these charts (even if flat), comment out or remove this block.
                        # actual_counts = counts_series.values
                        # all_top_n_counts_are_one = all(c == 1 for c in actual_counts)
                        
                        # if all_top_n_counts_are_one and len(actual_counts) >= 3:
                        #     logger.info(f"INFO: Generic chart for '{col_name_from_key}' in {report_type} would be flat; consider if suppression is desired.")
                        #     # If choosing to suppress:
                        #     # if 'insights' not in analysis: analysis['insights'] = {}
                        #     # analysis['insights'][f'suppressed_chart_info_{col_name_from_key.replace(" ", "_")}'] = (
                        #     #     f"A distribution chart for '{col_name_from_key}' was not generated because the top {len(actual_counts)} "
                        #     #     f"categories each had only one occurrence, indicating high uniqueness in this segment of data."
                        #     # )
                        #     # continue 
                        # ---- END: Reverted Suppression Logic ----

                        chart_bytes = ReportService._create_bar_chart(
                            data=counts_series,
                            title=chart_title,
                            xlabel=col_name_from_key.title(),
                            ylabel='# Records',
                            palette='Spectral',
                            figsize=(8, max(4, len(counts_series.index) * 0.5))
                        )
                        if chart_bytes:
                            charts[chart_key_name] = chart_bytes
                            chart_count += 1
                except Exception as e:
                    logger.error(f"Error generating generic chart for key '{key}' in {report_type}: {e}", exc_info=True)
        
        # This note is now more accurate: it appears if no charts were generated for any reason
        # (no suitable stats, or other limits hit, not just suppression).
        if chart_count == 0 and stats.get('record_count', 0) > 0 :
            logger.info(f"No generic charts were generated for {report_type} despite having data, due to data characteristics or limits.")
            if 'insights' not in analysis: analysis['insights'] = {}
            
            # Check if a more specific suppression insight for 'username' (or other columns) already exists
            # to avoid redundant "no chart" messages.
            # This is a simple check; a more robust way would be to track suppressed charts.
            has_suppression_insight = any(k.startswith('suppressed_chart_info_') for k in analysis.get('insights', {}).keys())

            if not has_suppression_insight: # Only add this general note if no specific suppression note was made
                analysis['insights']['no_generic_charts_note'] = (
                    "No generic distribution charts were generated for this section. "
                    "This can occur if data columns have too many unique values, too few unique values (e.g., only one distinct value after filtering), "
                    "or if other chart generation limits were met."
                )

        logger.debug(f"Generated {len(charts)} generic charts for {report_type}")
        return charts

    # --- Charting Helpers ---
    @staticmethod
    def _create_bar_chart(data: pd.Series, title: str, xlabel: str, ylabel: str, figsize=(8,4), palette="viridis", add_labels=True) -> Optional[BytesIO]:
       # (Keep existing implementation)
       if data.empty: return None
       try:
           fig, ax = ReportService._setup_chart(figsize=figsize); x_data = data.index.astype(str); y_data = data.values
           bars = sns.barplot(x=x_data, y=y_data, ax=ax, hue=x_data, palette=palette, width=0.6, legend=False )
           if add_labels:
                try: ReportService._add_bar_labels(ax, bars.patches, location='edge', fontsize=9)
                except AttributeError: ReportService._add_bar_labels(ax, bars, location='edge', fontsize=9)
           ax.set_title(title, fontsize=14, fontweight='bold'); ax.set_ylabel(ylabel, fontsize=12); ax.set_xlabel(xlabel, fontsize=12); ax.set_ylim(0, max(y_data.max() * 1.15, 1))
           if len(x_data) > 5 or max(len(label) for label in x_data) > 15: plt.xticks(rotation=45, ha='right')
           else: plt.xticks(rotation=0)
           plt.grid(axis='y', linestyle='--', alpha=0.6); return ReportService._save_plot_to_bytes(fig)
       except Exception as e: logger.error(f"Error generating bar chart '{title}': {e}", exc_info=True); return None
       finally: plt.close('all')

    @staticmethod
    def _create_pie_chart(data: pd.Series, title: str, figsize=(7, 5)) -> Optional[BytesIO]:
        # (Keep existing implementation)
        if data.empty: return None
        try:
             fig, ax = ReportService._setup_chart(figsize=figsize)
             explode = [0.1 if i == data.argmin() else 0 for i in range(len(data))] if len(data) > 3 else None
             wedges, texts, autotexts = ax.pie( data.values, labels=data.index.astype(str), autopct='%1.1f%%', startangle=90, pctdistance=0.85, explode=explode, colors=sns.color_palette("viridis", len(data)), wedgeprops={'edgecolor': 'white', 'linewidth': 1}, textprops={'fontsize': 9})
             for autotext in autotexts: autotext.set_color('white'); autotext.set_fontweight('bold')
             ax.set_title(title, fontsize=14, fontweight='bold', pad=20); ax.axis('equal')
             if len(data) > 6: plt.legend(wedges, data.index.astype(str), title="Categories", loc="center left", bbox_to_anchor=(1, 0, 0.5, 1)); [txt.set_visible(False) for txt in texts]; [autotext.set_visible(False) for autotext in autotexts]
             return ReportService._save_plot_to_bytes(fig)
        except Exception as e: logger.error(f"Error generating pie chart '{title}': {e}", exc_info=True); return None
        finally: plt.close('all')

    @staticmethod
    def _create_activity_heatmap(df: pd.DataFrame, date_col: str) -> Optional[BytesIO]:
       # (Keep existing implementation)
       if df.empty or date_col not in df.columns: return None
       try:
           if not pd.api.types.is_datetime64_any_dtype(df[date_col]): df[date_col] = pd.to_datetime(df[date_col], errors='coerce')
           if df[date_col].dt.tz is not None: df[date_col] = df[date_col].dt.tz_localize(None)
           df.dropna(subset=[date_col], inplace=True);
           if df.empty: return None
           df['hour'] = df[date_col].dt.hour; df['day_of_week'] = df[date_col].dt.day_name(); day_order = ['Monday', 'Tuesday', 'Wednesday', 'Thursday', 'Friday', 'Saturday', 'Sunday']
           hour_day = pd.crosstab(df['hour'], df['day_of_week']).reindex(index=range(24), columns=day_order, fill_value=0)
           if hour_day.sum().sum() == 0: return None
           fig, ax = ReportService._setup_chart(figsize=(10, 8))
           heatmap = sns.heatmap(hour_day, cmap="YlGnBu", linewidths=0.5, linecolor='lightgrey', ax=ax, cbar_kws={'label': 'Number of Submissions'}, annot=True, fmt="d", annot_kws={'fontsize': 8})
           hour_labels = {h: f"{h % 12 if h % 12 != 0 else 12} {'AM' if h < 12 else 'PM'}" for h in range(24)}
           ax.set_yticks(np.arange(len(hour_day.index)) + 0.5); ax.set_yticklabels([hour_labels.get(h, '') for h in hour_day.index], rotation=0)
           ax.set_title('Submission Activity by Hour and Day', fontsize=14, fontweight='bold'); ax.set_ylabel('Hour of Day', fontsize=12); ax.set_xlabel('Day of Week', fontsize=12); plt.xticks(rotation=30, ha='right')
           return ReportService._save_plot_to_bytes(fig)
       except Exception as e: logger.error(f"Error creating activity heatmap: {e}", exc_info=True); return None
       finally: plt.close('all')

    # --- Report Format Generation Helpers (Docs, PDF) ---
    @staticmethod
    def _add_document_title(doc_builder: Union[list, DocxDocumentClass, Presentation], title: str, level: int = 0):
       # (Keep existing implementation)
        try:
            if isinstance(doc_builder, list): styles = getSampleStyleSheet(); style_key = f'h{level+1}' if 0 < level < 6 else 'h1'; doc_builder.append(Paragraph(title, styles[style_key])); doc_builder.append(Spacer(1, 0.2*inch if level > 0 else 0.3*inch))
            elif isinstance(doc_builder, DocxDocumentClass): heading = doc_builder.add_heading(title, level=level); heading.alignment = WD_ALIGN_PARAGRAPH.CENTER if level == 0 else WD_ALIGN_PARAGRAPH.LEFT
            elif isinstance(doc_builder, Presentation): logger.debug("Title for PPTX handled in slide creation.")
            else: logger.warning(f"Unsupported doc type for title: {type(doc_builder)}")
        except Exception as e: logger.error(f"Error adding doc title '{title}': {e}")

    @staticmethod
    def _add_text_section(doc_builder: Union[list, DocxDocumentClass], heading: str, content: Dict[str, str], level: int = 2):
       # (Keep existing implementation)
       if not content: return
       styles = getSampleStyleSheet()
       try:
           if isinstance(doc_builder, list): style_key = f'h{level+1}' if 0 < level < 6 else 'h2'; doc_builder.append(Paragraph(heading, styles[style_key]))
           elif isinstance(doc_builder, DocxDocumentClass): doc_builder.add_heading(heading, level=level)
           is_stats = "statistic" in heading.lower()
           for key, text in content.items():
               text_str = str(text)
               if isinstance(doc_builder, list):
                   if is_stats: formatted_text = f"<b>{key.replace('_',' ').title()}:</b> {text_str}"; doc_builder.append(Paragraph(formatted_text, styles['Normal']))
                   else: doc_builder.append(Paragraph(f"• {text_str}", styles['Normal'], bulletText='•'))
               elif isinstance(doc_builder, DocxDocumentClass):
                   p = doc_builder.add_paragraph(style='List Bullet' if not is_stats else None)
                   if is_stats: p.add_run(f"{key.replace('_',' ').title()}: ").bold = True; p.add_run(text_str)
                   else: p.add_run(text_str)
           if isinstance(doc_builder, list): doc_builder.append(Spacer(1, 0.2*inch))
           elif isinstance(doc_builder, DocxDocumentClass): doc_builder.add_paragraph()
       except Exception as e: logger.error(f"Error adding text section '{heading}': {e}")

    @staticmethod
    def _add_chart_to_document(doc_builder: Union[list, DocxDocumentClass], chart_key: str, chart_bytes: Optional[BytesIO]):
        # (Keep existing implementation)
        if not isinstance(chart_bytes, BytesIO):
             logger.warning(f"Chart data for '{chart_key}' invalid."); error_text = f"[Chart '{chart_key.replace('_', ' ').title()}' could not be generated]"
             try:
                 if isinstance(doc_builder, list): styles = getSampleStyleSheet(); doc_builder.append(Paragraph(error_text, styles['Italic'])); doc_builder.append(Spacer(1, 0.3*inch))
                 elif isinstance(doc_builder, DocxDocumentClass): p = doc_builder.add_paragraph(); p.add_run(error_text).italic = True; p.alignment = WD_ALIGN_PARAGRAPH.CENTER; doc_builder.add_paragraph()
             except Exception: pass; return
        chart_bytes.seek(0); chart_title = chart_key.replace('_', ' ').title(); styles = getSampleStyleSheet()
        try:
            if isinstance(doc_builder, list): img = RLImage(chart_bytes, width=DEFAULT_CHART_WIDTH, height=DEFAULT_CHART_HEIGHT); img.hAlign = 'CENTER'; doc_builder.append(img); centered_caption_style = ParagraphStyle(name='CenteredCaption', parent=styles['Italic'], alignment=TA_CENTER); doc_builder.append(Paragraph(chart_title, centered_caption_style)); doc_builder.append(Spacer(1, 0.3*inch))
            elif isinstance(doc_builder, DocxDocumentClass): p_img = doc_builder.add_paragraph(); p_img.alignment = WD_ALIGN_PARAGRAPH.CENTER; run = p_img.add_run(); run.add_picture(chart_bytes, width=Inches(6.0)); p_caption = doc_builder.add_paragraph(style='Caption'); p_caption.add_run(chart_title); p_caption.alignment = WD_ALIGN_PARAGRAPH.CENTER; doc_builder.add_paragraph()
            else: logger.warning(f"Unsupported doc type for chart: {type(doc_builder)}")
        except Exception as e:
             logger.error(f"Error adding chart '{chart_key}': {e}", exc_info=True); error_text = f"[Error adding chart: {chart_key}]"
             try:
                 if isinstance(doc_builder, list): doc_builder.append(Paragraph(error_text, styles['Normal']))
                 elif isinstance(doc_builder, DocxDocumentClass): doc_builder.add_paragraph(error_text)
             except Exception: pass

    # --- Report Generation Orchestrator ---
    @staticmethod
    def _generate_report_data(report_params: dict, user: User) -> ProcessedData:
        """Fetches, flattens, and analyzes data for requested report types based on merged params."""
        report_type_req = report_params.get("report_type")
        processed_data: ProcessedData = {}

        if not report_type_req:
             processed_data['_error'] = {'error': "Missing 'report_type' in final parameters."}; return processed_data

        if report_type_req == "all": report_types_to_process = list(ReportService.ENTITY_CONFIG.keys())
        elif isinstance(report_type_req, list): report_types_to_process = report_type_req
        elif isinstance(report_type_req, str): report_types_to_process = [report_type_req]
        else: processed_data['_error'] = {'error': "Invalid report_type parameter."}; return processed_data

        # Determine if detailed parameters (columns, filters, sort) are present in the FINAL merged parameters
        has_detailed_params = any(k in report_params for k in ['columns', 'filters', 'sort_by'])
        is_single_entity_request = len(report_types_to_process) == 1

        question_info_map = {}
        if 'form_submissions' in report_types_to_process:
            try:
                questions_from_db = db.session.query(Question.text, QuestionType.type).join(QuestionType, Question.question_type_id == QuestionType.id).filter(Question.is_deleted == False).all()
                question_info_map = {q_text: q_type for q_text, q_type in questions_from_db}
            except Exception as q_err: logger.error(f"Could not pre-fetch question info: {q_err}")

        for report_type in report_types_to_process:
            if report_type not in ReportService.ENTITY_CONFIG: processed_data[report_type] = {'error': f"Unsupported report type: {report_type}"}; continue
            config = ReportService.ENTITY_CONFIG[report_type]; model_cls = config.get('model'); permission_entity = config.get('view_permission_entity')
            if not model_cls or not permission_entity: processed_data[report_type] = {'error': f"Configuration incomplete for: {report_type}"}; continue

            # Permission check for viewing the entity itself
            if not PermissionManager.has_permission(user, "view", permission_entity):
                processed_data[report_type] = {'error': f"Permission denied for {report_type} report."};
                logger.warning(f"User '{user.username}' permission denied for viewing entity '{permission_entity.value}'.")
                continue

            # Use detailed params ONLY if it's a single entity request AND detailed params were provided
            # Otherwise, use default columns logic.
            use_detailed_params = is_single_entity_request and has_detailed_params

            columns = None; filters = []; sort_by = config.get('default_sort', [])

            if use_detailed_params:
                # Use parameters from the final merged dictionary
                columns = report_params.get("columns") # Use user's columns if provided
                filters = report_params.get("filters", [])
                sort_by = report_params.get("sort_by", config.get('default_sort', []))
                logger.debug(f"Using detailed params for {report_type}: cols={len(columns) if columns else 'None'}, filters={len(filters)}, sort={len(sort_by)}")
            else:
                # --- Default Column Logic ---
                # This logic runs if it's a multi-entity request OR if no detailed params were provided
                logger.debug(f"Using default column logic for {report_type} (single_entity={is_single_entity_request}, has_detailed={has_detailed_params}).")
                try:
                    columns = list(sqla_inspect(model_cls).columns.keys())
                    logger.debug(f"Using default all direct columns for {report_type}: {len(columns)} columns")
                except Exception as inspect_err:
                     logger.error(f"Could not inspect columns for {model_cls.__name__}: {inspect_err}")
                     columns = config.get('default_columns', []) # Fallback to config default
                     if not columns:
                         processed_data[report_type] = {'error': f"Could not determine default columns for {report_type}."}; continue
                # Add default related columns if defined in config and not already present from inspection
                default_related_cols = config.get('default_columns', [])
                if default_related_cols:
                    inspected_cols_set = set(columns)
                    for d_col in default_related_cols:
                        if '.' in d_col and d_col not in inspected_cols_set:
                            columns.append(d_col)
                # Reset filters and sort to defaults when not using detailed params
                filters = []
                sort_by = config.get('default_sort', [])


            if not columns or not isinstance(columns, list):
                logger.error(f"Column list for {report_type} ended up invalid: {columns}")
                processed_data[report_type] = {'error': f"Could not determine columns for {report_type}."}
                continue

            # --- Filter Sensitive Columns ---
            if report_type == 'users' and not (user.role and user.role.is_super_user):
                sensitive_columns = ['password_hash']
                original_column_count = len(columns)
                columns = [col for col in columns if col not in sensitive_columns]
                if len(columns) < original_column_count: logger.info(f"Removed sensitive columns for non-admin user '{user.username}'.")

            if not columns: processed_data[report_type] = {'error': f"No accessible columns for {report_type}."}; continue

            # Construct parameters specific to this entity's processing
            # Use values from the final merged report_params
            current_entity_params = {
                "columns": columns,
                "filters": filters,
                "sort_by": sort_by,
                "report_type": report_type, # Specific entity type
                "report_title": report_params.get("report_title", DEFAULT_REPORT_TITLE), # Global title
                "output_format": report_params.get("output_format", "xlsx").lower(), # Global format
                "_internal_question_info": question_info_map if report_type == 'form_submissions' else {},
                "_internal_config": config,
                "sheet_name": report_params.get(f"{report_type}_sheet_name", report_type.replace("_", " ").title()), # Sheet name can be overridden
                "table_options": report_params.get(f"{report_type}_table_options", report_params.get("table_options", {})), # Table options can be overridden
                "include_data_table_in_ppt": report_params.get("include_data_table_in_ppt", False), # Global PPT option
                "max_ppt_table_rows": report_params.get("max_ppt_table_rows", DEFAULT_PPTX_TABLE_ROWS), # Global PPT option
            }
            try:
                logger.info(f"Processing report type '{report_type}' for user '{user.username}'...")
                fetched_objects = ReportService._fetch_data(model_cls, filters, sort_by, user, columns)
                data = ReportService._flatten_data(fetched_objects, columns, report_type)
                # Pass the entity-specific params to analysis
                analysis_results = ReportService._analyze_data(data, current_entity_params, report_type)
                processed_data[report_type] = {'error': None, 'data': data, 'objects': fetched_objects, 'params': current_entity_params, 'analysis': analysis_results }
                logger.info(f"Successfully processed '{report_type}' ({len(data)} records) for user '{user.username}'.")
            except Exception as proc_err:
                logger.error(f"Error processing data for {report_type} for user '{user.username}': {proc_err}", exc_info=True)
                processed_data[report_type] = {'error': f"Error processing data: {proc_err}", 'params': current_entity_params, 'analysis': {}, 'data': [], 'objects':[]}
        return processed_data

    # --- Format-Specific Generation Functions (XLSX, CSV, PDF, DOCX - Unchanged) ---
    @staticmethod
    def _generate_xlsx_report(processed_data: ProcessedData, global_params: ParamsDict) -> BytesIO:
        # (Keep existing implementation with ID sort)
        output = BytesIO()
        with xlsxwriter.Workbook(output, {'in_memory': True, 'remove_timezone': True, 'strings_to_numbers': False, 'strings_to_formulas': False}) as workbook:
            header_format = workbook.add_format({'bold': True, 'text_wrap': True, 'valign': 'top', 'fg_color': '#D7E4BC', 'border': 1, 'align': 'center'})
            wrap_format = workbook.add_format({'text_wrap': True, 'valign': 'top', 'align': 'left'})

            for report_type, result in processed_data.items():
                if result.get('error'):
                    try: sheet_name = f"ERROR_{report_type}"[:MAX_XLSX_SHEET_NAME_LEN]; worksheet = workbook.add_worksheet(sheet_name); worksheet.write(0, 0, f"Error: {report_type}:"); worksheet.write(1, 0, str(result['error'])); worksheet.set_column(0, 0, 100, wrap_format)
                    except Exception as sheet_err: logger.error(f"Could not write error sheet for {report_type}: {sheet_err}")
                    continue

                params = result.get('params', {}); data = result.get('data', []); columns = params.get('columns', []); analysis = result.get('analysis', {}); sheet_name = params.get("sheet_name", report_type.replace("_", " ").title())[:MAX_XLSX_SHEET_NAME_LEN]; table_options = params.get('table_options', {})

                # Sort data list by 'id' if present
                if data and columns and 'id' in columns:
                    try:
                        # Attempt numeric sort first, placing non-numeric/missing IDs last
                        data.sort(key=lambda row: int(row['id']) if isinstance(row.get('id'), (int, float)) or (isinstance(row.get('id'), str) and row.get('id', '').isdigit()) else float('inf'))
                        logger.debug(f"Sorted data for sheet '{sheet_name}' by ID ascending.")
                    except (ValueError, TypeError) as sort_err:
                         logger.warning(f"Could not sort sheet '{sheet_name}' numerically by ID, attempting string sort: {sort_err}")
                         try: data.sort(key=lambda row: str(row.get('id', ''))) # Fallback string sort
                         except Exception as sort_err_str: logger.error(f"String sort by ID also failed for sheet '{sheet_name}': {sort_err_str}")
                    except Exception as sort_err: logger.error(f"Error sorting data for sheet '{sheet_name}' by ID: {sort_err}", exc_info=True)

                if not columns: logger.warning(f"Skipping sheet {report_type}: No columns."); continue

                try:
                    worksheet = workbook.add_worksheet(sheet_name); current_row = 0
                    worksheet.merge_range(current_row, 0, current_row, max(3, len(columns)-1), f"Report: {sheet_name}", workbook.add_format({'bold': True, 'font_size': 14, 'align': 'center'})); current_row += 1
                    worksheet.merge_range(current_row, 0, current_row, max(3, len(columns)-1), f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", workbook.add_format({'italic': True, 'align': 'center'})); current_row += 2

                    table_data = []; col_max_lens = {i: len(str(columns[i])) for i in range(len(columns))}
                    for row_dict in data: # Iterate over potentially sorted data
                        row_values = []
                        for c_idx, col_key in enumerate(columns):
                            cell_value = row_dict.get(col_key)
                            if cell_value is None: formatted_value = ''
                            elif isinstance(cell_value, bool): formatted_value = "Yes" if cell_value else "No"
                            elif isinstance(cell_value, (datetime, date)):
                                try: formatted_value = cell_value; col_max_lens[c_idx] = max(col_max_lens.get(c_idx, 0), 20) # Fixed width for dates
                                except ValueError: formatted_value = str(cell_value.isoformat()); col_max_lens[c_idx] = max(col_max_lens.get(c_idx, 0), len(formatted_value))
                            else: formatted_value = str(cell_value); col_max_lens[c_idx] = max(col_max_lens.get(c_idx, 0), len(formatted_value))
                            row_values.append(formatted_value)
                        table_data.append(row_values)

                    table_headers = [{'header': col.replace(".", " ").replace("_", " ").title()} for col in columns]
                    first_row_table = current_row; first_col_table = 0

                    if not data and not columns: worksheet.write(first_row_table, 0, "No data or columns."); current_row +=1
                    elif not data:
                         for col_idx, header_info in enumerate(table_headers): worksheet.write(first_row_table, col_idx, header_info['header'], header_format)
                         current_row +=1
                    else:
                         last_row_table = first_row_table + len(table_data); last_col_table = first_col_table + len(columns) - 1
                         worksheet.add_table(first_row_table, first_col_table, last_row_table, last_col_table, {'data': table_data, 'columns': table_headers, 'style': table_options.get('style', 'Table Style Medium 9'), 'name': f"{sheet_name.replace(' ','_')}_Table", 'header_row': True, 'banded_rows': table_options.get('banded_rows', True), 'autofilter': table_options.get('autofilter', True)})
                         current_row = last_row_table + 1

                    for col_idx, max_len in col_max_lens.items(): width = min(max(max_len, 10) + 2, 60); worksheet.set_column(col_idx, col_idx, width, wrap_format)
                    chart_start_row = current_row + 2; chart_col = 1
                    for chart_name, chart_bytes in analysis.get('charts', {}).items():
                        if isinstance(chart_bytes, BytesIO):
                            try: chart_bytes.seek(0); worksheet.write(chart_start_row, chart_col - 1, f"{chart_name.replace('_',' ').title()}:"); worksheet.insert_image(chart_start_row + 1, chart_col, f"chart_{chart_name}.png", {'image_data': chart_bytes, 'x_scale': 0.6, 'y_scale': 0.6}); chart_start_row += 20
                            except Exception as chart_err: logger.error(f"Failed chart insert {chart_name}: {chart_err}"); worksheet.write(chart_start_row, chart_col - 1, f"Error chart: {chart_name}"); chart_start_row += 1
                except Exception as sheet_err:
                    logger.error(f"Failed sheet '{sheet_name}': {sheet_err}", exc_info=True)
                    try:
                        error_sheet_name = f"ERROR_{sheet_name[:25]}"
                        if error_sheet_name not in workbook.sheetnames:
                            error_worksheet = workbook.add_worksheet(error_sheet_name)
                            error_worksheet.write(0, 0, f"Failed generating original sheet '{sheet_name}'. Error:")
                            error_msg_str = str(sheet_err)
                            max_cell_len = 32767
                            error_worksheet.write(1, 0, error_msg_str[:max_cell_len])
                            error_worksheet.set_column(0, 0, 100)
                            logger.info(f"Created error sheet '{error_sheet_name}' due to failure on '{sheet_name}'.")
                        else: logger.warning(f"Attempted to create error sheet '{error_sheet_name}', but it already exists.")
                    except Exception as inner_err: logger.error(f"Could not write error sheet '{error_sheet_name}' after initial failure on '{sheet_name}': {inner_err}", exc_info=True); pass
        output.seek(0); return output

    @staticmethod
    def _generate_csv_report(processed_data: ProcessedData, global_params: ParamsDict) -> BytesIO:
        # (Keep existing implementation)
        if len(processed_data) == 1:
            report_type = list(processed_data.keys())[0]; result = processed_data[report_type]
            if result.get('error'): raise ValueError(f"CSV Error: {result['error']}")
            data = result.get('data', []); columns = result.get('params', {}).get('columns', [])
            if not columns: raise ValueError("No columns for CSV.")
            output = StringIO(); writer = csv.DictWriter(output, fieldnames=columns, quoting=csv.QUOTE_MINIMAL, extrasaction='ignore'); writer.writeheader()
            formatted_data = []
            for row_dict in data:
                 new_row = {}
                 for col in columns: val = row_dict.get(col); new_row[col] = val.isoformat() if isinstance(val, (datetime, date)) else ('' if val is None else str(val))
                 formatted_data.append(new_row)
            writer.writerows(formatted_data); output.seek(0); return BytesIO(output.getvalue().encode('utf-8'))
        else:
            zip_buffer = BytesIO()
            with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zipf:
                for report_type, result in processed_data.items():
                    filename_base = result.get('params', {}).get('sheet_name', report_type)
                    if result.get('error'): zipf.writestr(f"{filename_base}_error.txt", f"Error: {result['error']}"); continue
                    data = result.get('data', []); columns = result.get('params', {}).get('columns', [])
                    if not columns: zipf.writestr(f"{filename_base}_error.txt", f"No columns."); continue
                    csv_output = StringIO(); writer = csv.DictWriter(csv_output, fieldnames=columns, quoting=csv.QUOTE_MINIMAL, extrasaction='ignore'); writer.writeheader()
                    formatted_data = []
                    for row_dict in data:
                        new_row = {}
                        for col in columns: val = row_dict.get(col); new_row[col] = val.isoformat() if isinstance(val, (datetime, date)) else ('' if val is None else str(val))
                        formatted_data.append(new_row)
                    writer.writerows(formatted_data); csv_output.seek(0)
                    zipf.writestr(f"{filename_base}.csv", csv_output.getvalue().encode('utf-8'))
            zip_buffer.seek(0); return zip_buffer

    @staticmethod
    def _generate_pdf_report(processed_data: ProcessedData, global_params: ParamsDict) -> BytesIO:
        # (Keep existing implementation)
        buffer = BytesIO(); doc = SimpleDocTemplate(buffer, pagesize=letter, leftMargin=0.75*inch, rightMargin=0.75*inch, topMargin=0.75*inch, bottomMargin=0.75*inch); styles = getSampleStyleSheet(); story = []
        ReportService._add_document_title(story, global_params.get("report_title", DEFAULT_REPORT_TITLE), level=0)
        first_section = True
        for report_type, result in processed_data.items():
            if result.get('error'): story.append(Paragraph(f"Error: {report_type}: {result['error']}", styles['Italic'])); story.append(Spacer(1, 0.2*inch)); continue
            analysis = result.get('analysis', {}); params = result.get('params', {}); section_title = params.get("sheet_name", report_type.replace("_", " ").title())
            if not first_section: story.append(PageBreak())
            first_section = False
            ReportService._add_document_title(story, section_title, level=1)
            if analysis.get('insights'): ReportService._add_text_section(story, "Key Insights:", analysis['insights'], level=2)
            if analysis.get('summary_stats'): simple_stats = {k:v for k,v in analysis['summary_stats'].items() if not k.startswith('_') and not isinstance(v, (dict, list, pd.DataFrame))}; ReportService._add_text_section(story, "Summary Statistics:", simple_stats, level=2)
            if analysis.get('charts'):
                 story.append(Paragraph("Visual Analysis:", styles['h3'])); story.append(Spacer(1, 0.1*inch))
                 for chart_key, chart_bytes in analysis['charts'].items(): ReportService._add_chart_to_document(story, chart_key, chart_bytes)
        try: doc.build(story)
        except Exception as build_err:
            logger.error(f"Error building PDF: {build_err}", exc_info=True); buffer = BytesIO(); doc = SimpleDocTemplate(buffer, pagesize=letter); story = [Paragraph("Error Building PDF Report", styles['h1']), Paragraph(str(build_err), styles['Normal'])];
            try: doc.build(story)
            except: return BytesIO(b"Failed PDF error report.")
        buffer.seek(0); return buffer

    @staticmethod
    def _generate_docx_report(processed_data: ProcessedData, global_params: ParamsDict) -> BytesIO:
        # (Keep existing implementation)
        document = docx.Document(); buffer = BytesIO()
        ReportService._add_document_title(document, global_params.get("report_title", DEFAULT_REPORT_TITLE), level=0)
        first_section = True
        for report_type, result in processed_data.items():
             if result.get('error'): document.add_paragraph(f"Error: {report_type}:").italic = True; document.add_paragraph(str(result['error'])).italic = True; document.add_paragraph(); continue
             analysis = result.get('analysis', {}); params = result.get('params', {}); section_title = params.get("sheet_name", report_type.replace("_", " ").title())
             if not first_section: document.add_page_break()
             first_section = False
             ReportService._add_document_title(document, section_title, level=1)
             if analysis.get('insights'): ReportService._add_text_section(document, "Key Insights", analysis['insights'], level=2)
             if analysis.get('summary_stats'): simple_stats = {k:v for k,v in analysis['summary_stats'].items() if not k.startswith('_') and not isinstance(v, (dict, list, pd.DataFrame))}; ReportService._add_text_section(document, "Summary Statistics", simple_stats, level=2)
             if analysis.get('charts'):
                 document.add_heading("Visual Analysis", level=2)
                 for chart_key, chart_bytes in analysis['charts'].items(): ReportService._add_chart_to_document(document, chart_key, chart_bytes)
        try: document.save(buffer)
        except Exception as save_err:
             logger.error(f"Error saving DOCX: {save_err}", exc_info=True); buffer = BytesIO(); doc_err = docx.Document(); doc_err.add_heading("Error Saving DOCX Report", level=0); doc_err.add_paragraph(str(save_err));
             try: doc_err.save(buffer)
             except: return BytesIO(b"Failed DOCX error report.")
        buffer.seek(0); return buffer

    # --- PPTX Specific Helpers ---
    @staticmethod
    def _add_pptx_text_slide(prs: Presentation, title: str, content: List[Tuple[str, Optional[int], Optional[bool]]]):
        # (Keep existing implementation)
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[1]); slide.shapes.title.text = title
            if len(slide.placeholders) > 1:
                tf = slide.placeholders[1].text_frame; tf.clear() ; tf.word_wrap = True
                for text, level, is_bold in content: p = tf.add_paragraph(); p.text = str(text); p.font.bold = is_bold or False; effective_level = min(level, 5) if level is not None else 0; p.level = effective_level; font_size = PptxPt(max(10, 18 - (effective_level * 2.5))); p.font.size = font_size
            else: logger.warning(f"Slide layout for '{title}' missing content placeholder.")
        except Exception as e: logger.error(f"Error adding PPTX text slide '{title}': {e}", exc_info=True)

    @staticmethod
    def _add_pptx_chart_slide(prs: Presentation, chart_key: str, chart_bytes: Optional[BytesIO], analysis: AnalysisDict):
        # (Keep existing implementation)
        if not isinstance(chart_bytes, BytesIO): logger.warning(f"Skipping PPTX chart '{chart_key}': Invalid data."); return
        chart_bytes.seek(0); chart_title = chart_key.replace('_', ' ').title()
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[5]); slide.shapes.title.text = chart_title
            pic_left, pic_top = PptxInches(0.5), PptxInches(1.25); pic_max_width, pic_max_height = PptxInches(9.0), PptxInches(5.0)
            pic = slide.shapes.add_picture(chart_bytes, pic_left, pic_top, width=pic_max_width)
            if pic.width > pic_max_width: ratio = float(pic_max_width) / float(pic.width); pic.width = int(pic_max_width); pic.height = int(float(pic.height) * ratio)
            if pic.height > pic_max_height: ratio = float(pic_max_height) / float(pic.height); pic.height = int(pic_max_height); pic.width = int(float(pic.width) * ratio)
            pic.left = int((prs.slide_width - pic.width) / 2); pic.top = PptxInches(1.25)
        except Exception as img_err: logger.error(f"Error adding chart '{chart_key}' to PPTX: {img_err}", exc_info=True)

    @staticmethod
    def _add_pptx_data_table_slide(prs: Presentation, title: str, data: DataList, params: ParamsDict):
        # (Keep existing implementation with Dark Blue Header Fix)
        max_rows = params.get("max_ppt_table_rows", DEFAULT_PPTX_TABLE_ROWS); columns = params.get("columns", []); data_sample = data[:max_rows]
        if not data_sample or not columns: return
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[5]); slide.shapes.title.text = title
            display_cols = [c for c in ["id", "name", "title", "submitted_by", "created_at", "updated_at", "status"] if c in columns]
            if not display_cols: display_cols = [c for c in columns if not c.startswith(ANSWERS_PREFIX)][:5]
            if not display_cols: display_cols = columns[:5]
            rows = len(data_sample) + 1; cols = len(display_cols); left, top = PptxInches(0.5), PptxInches(1.5); width = PptxInches(min(9.0, cols * 2.0)); height = PptxInches(min(5.5, rows * 0.35))
            shape = slide.shapes.add_table(rows, cols, left, top, width, height); table = shape.table
            for c_idx, col_name in enumerate(display_cols):
                cell = table.cell(0, c_idx); clean_col_name = col_name.split(ANSWERS_PREFIX, 1)[-1] if col_name.startswith(ANSWERS_PREFIX) else col_name; cell.text = clean_col_name.replace(".", " ").replace("_", " ").title()[:25]
                p = cell.text_frame.paragraphs[0]; p.font.bold = True; p.font.size = PptxPt(10); p.alignment = PP_ALIGN.CENTER; p.font.color.rgb = PptxRGBColor(255, 255, 255) # White text
                cell.fill.solid(); cell.fill.fore_color.rgb = PptxRGBColor(0, 0, 128); cell.vertical_anchor = MSO_ANCHOR.MIDDLE # Dark Blue BG
            for r_idx, row_dict in enumerate(data_sample):
                for c_idx, col_name in enumerate(display_cols):
                    value = row_dict.get(col_name, ''); str_value = str(value);
                    if len(str_value) > 75: str_value = str_value[:72] + "..."
                    cell = table.cell(r_idx + 1, c_idx); cell.text = str_value; p = cell.text_frame.paragraphs[0]; p.font.size = PptxPt(9); cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if len(data) > max_rows:
                note_top = shape.top + shape.height + PptxInches(0.1)
                if note_top < prs.slide_height - PptxInches(0.5): note_shape = slide.shapes.add_textbox(left, note_top, width, PptxInches(0.5)); p_note = note_shape.text_frame.paragraphs[0]; p_note.text = f"Note: Showing first {max_rows} of {len(data)} records."; p_note.font.size = PptxPt(10); p_note.font.italic = True; p_note.alignment = PP_ALIGN.CENTER
        except Exception as table_err: logger.error(f"Error creating PPTX table slide '{title}': {table_err}", exc_info=True)

    # --- Specific and Default PPTX Generators ---
    @staticmethod
    def _generate_submission_pptx(result: Dict[str, Any], global_params: ParamsDict) -> BytesIO:
        # (Keep existing implementation)
        prs = Presentation(); buffer = BytesIO(); analysis = result.get('analysis', {}); params = result.get('params', {}); data = result.get('data', [])
        report_title = params.get("report_title", "Submission Analysis"); report_type_str = params.get("report_type", "").replace('_',' ').title()
        slide_title = prs.slides.add_slide(prs.slide_layouts[0]); slide_title.shapes.title.text = report_title
        try: slide_title.placeholders[1].text = f"Report Type: {report_type_str}\nGenerated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        except IndexError: pass
        stats = analysis.get('summary_stats', {}); insights = analysis.get('insights', {})
        summary_content: List[Tuple[str, Optional[int], Optional[bool]]] = []
        if stats.get('record_count'): summary_content.append((f"Total Submissions: {stats['record_count']}", 0, True))
        if stats.get('overall_submission_range'): summary_content.append((f"Date Range: {stats['overall_submission_range']['first'].split(' ')[0]} to {stats['overall_submission_range']['last'].split(' ')[0]}", 1, False))
        if stats.get('average_daily_submissions'): summary_content.append((f"Avg Daily Rate: {stats['average_daily_submissions']:.1f}", 1, False))
        if insights: summary_content.append(("", None, False)); summary_content.append(("Key Insights:", 0, True)); [summary_content.append((f"{insight_text}", 1, False)) for key, insight_text in insights.items() if key != 'status']
        if summary_content: ReportService._add_pptx_text_slide(prs, "Executive Summary", summary_content)
        for chart_key, chart_bytes in analysis.get('charts', {}).items(): ReportService._add_pptx_chart_slide(prs, chart_key, chart_bytes, analysis)
        if data and params.get("include_data_table_in_ppt", False): ReportService._add_pptx_data_table_slide(prs, "Submission Data Sample", data, params)
        ReportService._add_pptx_text_slide(prs, "Conclusions", [ ("Review completed.", 0, True), (insights.get("activity_rate", ""), 1, False), (insights.get("top_department", ""), 1, False), (insights.get("top_user", ""), 1, False), ])
        prs.save(buffer); buffer.seek(0); return buffer

    @staticmethod
    def _generate_default_pptx(result: Dict[str, Any], global_params: ParamsDict) -> BytesIO:
        # (Keep existing implementation)
        prs = Presentation(); buffer = BytesIO(); analysis = result.get('analysis', {}); params = result.get('params', {}); data = result.get('data', [])
        report_title = params.get("report_title", "Data Analysis"); report_type_str = params.get("report_type", "Unknown").replace('_',' ').title(); entity_title = params.get("sheet_name", report_type_str)
        logger.info(f"Generating default PPTX for: {report_type_str}")
        slide_title = prs.slides.add_slide(prs.slide_layouts[0]) ; slide_title.shapes.title.text = report_title
        try: subtitle = slide_title.placeholders[1]; subtitle.text = f"Entity: {entity_title}\nGenerated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        except IndexError: pass
        stats = analysis.get('summary_stats', {}); insights = analysis.get('insights', {})
        summary_content: List[Tuple[str, Optional[int], Optional[bool]]] = []
        if stats.get('record_count') is not None: summary_content.append((f"Total Records: {stats['record_count']}", 0, True))
        simple_stats = {k:v for k,v in stats.items() if not k.startswith('_') and k != 'record_count' and not isinstance(v, (dict, list, pd.DataFrame))}
        if simple_stats: summary_content.append(("", None, False)); summary_content.append(("Key Statistics:", 0, True)); [summary_content.append((f"{key.replace('_',' ').title()}: {stat_value}", 1, False)) for key, stat_value in simple_stats.items()]
        cleaned_insights = {k: v for k, v in insights.items() if k != 'status'}
        if cleaned_insights: summary_content.append(("", None, False)); summary_content.append(("Key Insights:", 0, True)); [summary_content.append((f"{insight_text}", 1, False)) for key, insight_text in cleaned_insights.items()]
        if summary_content: ReportService._add_pptx_text_slide(prs, f"{entity_title} - Summary", summary_content)
        else: ReportService._add_pptx_text_slide(prs, f"{entity_title} - Summary", [("No specific stats/insights generated.", 0, False)])
        charts = analysis.get('charts', {})
        if charts: logger.info(f"Adding {len(charts)} charts for {report_type_str}"); [ReportService._add_pptx_chart_slide(prs, ck, cb, analysis) for ck, cb in charts.items()]
        else: logger.info(f"No charts for {report_type_str}")
        if data and params.get("include_data_table_in_ppt", False): logger.info(f"Adding data table for {report_type_str}"); ReportService._add_pptx_data_table_slide(prs, f"{entity_title} - Data Sample", data, params)
        ReportService._add_pptx_text_slide(prs, "Report End", [ (f"Analysis for {entity_title} completed.", 0, True), (f"Total records: {stats.get('record_count', 'N/A')}.", 1, False) ])
        try: prs.save(buffer); buffer.seek(0); logger.info(f"Saved default PPTX buffer for {report_type_str}"); return buffer
        except Exception as e: logger.error(f"Failed to save default PPTX: {e}", exc_info=True); raise

    # --- PPTX Router ---
    @staticmethod
    def _generate_pptx_report(processed_data: ProcessedData, global_params: ParamsDict) -> BytesIO:
        # (Keep existing implementation - uses default fallback)
        primary_report_type = None; primary_result = None
        for rt, res in processed_data.items():
            if not res.get('error'): primary_report_type = rt; primary_result = res; logger.info(f"PPTX focus entity: {primary_report_type}"); break
        if not primary_report_type or not primary_result:
            all_errors = "; ".join([f"{rt}: {res['error']}" for rt, res in processed_data.items() if res.get('error')]); error_msg = f"No valid data for PPTX. Errors: {all_errors}" if all_errors else "No valid data for PPTX."; logger.error(error_msg)
            prs = Presentation(); ReportService._add_pptx_text_slide(prs, "Report Error", [(error_msg, 0, False)]); buffer = BytesIO()
            try: prs.save(buffer); buffer.seek(0); return buffer
            except Exception as save_err: logger.error(f"Failed to save error PPTX: {save_err}"); raise RuntimeError("Failed report & error report.")
        params = primary_result.get('params', {}); config = params.get('_internal_config', {}); pptx_gen_func_name = config.get('format_generators', {}).get('pptx')
        if pptx_gen_func_name and hasattr(ReportService, pptx_gen_func_name):
            try: generator_func: PptxGenerator = getattr(ReportService, pptx_gen_func_name); logger.info(f"Using specific PPTX gen '{pptx_gen_func_name}'."); return generator_func(primary_result, global_params)
            except Exception as specific_gen_err: logger.error(f"Error in specific PPTX gen '{pptx_gen_func_name}': {specific_gen_err}", exc_info=True); logger.warning(f"Falling back to default PPTX gen.")
        logger.info(f"Using default PPTX gen for {primary_report_type}.")
        try: return ReportService._generate_default_pptx(primary_result, global_params)
        except Exception as default_gen_err:
            logger.error(f"Error in default PPTX gen for {primary_report_type}: {default_gen_err}", exc_info=True)
            prs = Presentation(); error_msg = f"Failed PPTX gen for {primary_report_type}: {default_gen_err}"; ReportService._add_pptx_text_slide(prs, "Report Error", [(error_msg, 0, False)]); buffer = BytesIO()
            try: prs.save(buffer); buffer.seek(0); return buffer
            except Exception as save_err: logger.error(f"Failed save error PPTX: {save_err}"); raise RuntimeError(f"Failed PPTX for {primary_report_type} & error report.") from default_gen_err

    # --- Main Public Method (Unchanged implementation from previous version) ---
    @staticmethod
    def generate_report(report_params: dict, user: User) -> ReportResult:
        """
        Generates a report based on parameters, handling permissions and multiple formats.
        Can load configuration from a saved ReportTemplate if 'template_id' is provided.

        Args:
            report_params (dict): Parameters defining the report. Can include 'template_id'.
                                  Direct parameters override template configuration.
            user (User): The user requesting the report.

        Returns:
            ReportResult: Tuple containing (BytesIO buffer | None, filename | None, mime_type | None, error_message | None)
        """
        final_report_params = {}
        template_id = report_params.get('template_id')
        request_params = report_params.copy() # Keep original request params separate

        # 1. Load Template Configuration (if template_id provided)
        if template_id:
            logger.info(f"Report request includes template_id: {template_id}. Loading template...")
            try:
                template_id = int(template_id)
                template = ReportTemplate.find_by_id(template_id)

                if not template:
                    logger.warning(f"Template ID {template_id} not found for report generation.")
                    return None, None, None, f"Report template with ID {template_id} not found."

                # Check permissions for the template
                is_owner = template.user_id == user.id
                is_admin = user.role and user.role.name == RoleType.ADMIN
                if not is_owner and not template.is_public and not is_admin:
                    logger.warning(f"User '{user.username}' permission denied for template ID {template_id}.")
                    return None, None, None, f"Permission denied to use report template ID {template_id}."

                # Load configuration from the template
                template_config = template.configuration
                if not isinstance(template_config, dict):
                     logger.error(f"Template ID {template_id} has invalid configuration (not a dict).")
                     return None, None, None, f"Report template {template_id} has invalid configuration."

                logger.info(f"Successfully loaded configuration from template '{template.name}' (ID: {template_id}).")
                final_report_params = template_config.copy() # Start with template config

            except ValueError:
                 logger.warning(f"Invalid template_id format: {template_id}")
                 return None, None, None, "Invalid template_id format. Must be an integer."
            except Exception as e:
                 logger.exception(f"Error loading report template ID {template_id}: {e}")
                 return None, None, None, f"An internal error occurred while loading the report template: {str(e)}"

        # 2. Merge/Override with Request Parameters
        # Parameters directly in the request override those from the template
        # Remove 'template_id' from request_params before merging if it exists
        request_params.pop('template_id', None)
        final_report_params.update(request_params) # Direct request params take precedence

        # --- Proceed with existing report generation logic using final_report_params ---
        try:
            # 3. Initial setup (format, filename base from final params)
            output_format = final_report_params.get("output_format", "xlsx").lower()
            if output_format not in SUPPORTED_FORMATS:
                 return None, None, None, f"Unsupported format: {output_format}. Supported formats: {', '.join(SUPPORTED_FORMATS)}"

            base_filename = final_report_params.get("filename") # Filename can come from template or request

            # 4. Process data for requested entities using final_report_params
            # The _generate_report_data method now receives the potentially merged params
            processed_data = ReportService._generate_report_data(final_report_params, user)

            # 5. Handle errors from data processing
            if '_error' in processed_data:
                return None, None, None, processed_data['_error']['error']
            if not any(not res.get('error') for res in processed_data.values()):
                all_errors = "; ".join([f"{rt}: {res['error']}" for rt, res in processed_data.items() if res.get('error')])
                error_msg = f"No data generated. Errors: {all_errors}" if all_errors else "No data found for the specified parameters."
                return None, None, None, error_msg

            # 6. Determine final filename (if not provided in params)
            report_type_req = final_report_params.get("report_type") # Get report type from final params
            if not base_filename:
                ts = datetime.now().strftime('%Y%m%d_%H%M'); name_part = "custom_report"
                if report_type_req == "all": name_part = "full_report"
                elif isinstance(report_type_req, list): name_part = "multi_report"
                elif isinstance(report_type_req, str): name_part = f"report_{report_type_req}"
                # If loaded from template, maybe use template name?
                if template_id and 'template' in locals() and template:
                    safe_template_name = "".join(c if c.isalnum() else "_" for c in template.name)
                    base_filename = f"template_{safe_template_name}_{ts}"
                else:
                     base_filename = f"{name_part}_{ts}"

            # 7. Prepare for format generation
            final_buffer: Optional[BytesIO] = None
            final_filename: str = f"{base_filename}.{output_format}"
            mime_type: Optional[str] = None
            generator_map: Dict[str, Callable] = { "xlsx": ReportService._generate_xlsx_report, "csv": ReportService._generate_csv_report, "pdf": ReportService._generate_pdf_report, "docx": ReportService._generate_docx_report, "pptx": ReportService._generate_pptx_report, }
            mime_map = { "xlsx": 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', "csv": 'application/zip' if isinstance(final_report_params.get("report_type"), list) and len(final_report_params.get("report_type",[])) > 1 else 'text/csv', "pdf": 'application/pdf', "docx": 'application/vnd.openxmlformats-officedocument.wordprocessingml.document', "pptx": 'application/vnd.openxmlformats-officedocument.presentationml.presentation', }

            # 8. Generate the report file using the appropriate generator
            if output_format in generator_map:
                try:
                    generator_func = generator_map[output_format]
                    # Pass final_report_params to the generator function if it needs them
                    # Currently, generators mainly use processed_data, but pass global params just in case
                    final_buffer = generator_func(processed_data, final_report_params)
                    mime_type = mime_map[output_format]
                    # Adjust filename for multi-CSV zip
                    if output_format == 'csv' and isinstance(final_report_params.get("report_type"), list) and len(final_report_params.get("report_type",[])) > 1:
                         final_filename = f"{base_filename}.zip"
                         mime_type = 'application/zip' # Ensure mime type is zip for multi-csv
                except Exception as format_err:
                    # Handle errors during format generation
                    logger.error(f"Error generating {output_format}: {format_err}", exc_info=True)
                    return None, None, None, f"Error during {output_format} generation: {format_err}"
            else:
                # This case should be caught earlier by the format check
                 return None, None, None, f"Unsupported format: {output_format}"

            # 9. Return result or handle buffer failure
            if final_buffer and isinstance(final_buffer, BytesIO):
                logger.info(f"{output_format.upper()} report '{final_filename}' generated successfully for user '{user.username}'. Template ID used: {template_id if template_id else 'None'}.")
                return final_buffer, final_filename, mime_type, None
            else:
                # This case means the generator function ran without error but didn't return a valid BytesIO buffer
                error_msg = f"Failed to generate buffer for {output_format}."
                logger.error(error_msg)
                return None, None, None, error_msg

        # Handle any unexpected errors during the whole process
        except Exception as e:
            logger.exception(f"Unexpected error in report generation for user '{user.username}', params: {final_report_params}: {e}")
            return None, None, None, f"An unexpected error occurred during report generation."

    # --- Schema Retrieval (Unchanged implementation from previous version) ---
    @staticmethod
    def get_database_schema() -> SchemaResult:
        """Retrieves database schema and table row counts."""
        # (Keep existing implementation)
        try:
            inspector = sqla_inspect(db.engine); schema_data = {}; table_names = inspector.get_table_names()
            for table_name in table_names:
                try:
                    cols = inspector.get_columns(table_name); pk_constraint = inspector.get_pk_constraint(table_name); pk_cols = pk_constraint.get('constrained_columns', []) if pk_constraint else []; fks = inspector.get_foreign_keys(table_name)
                    columns_info = [{"name": c['name'], "type": str(c['type']), "nullable": c.get('nullable', True), "default": str(c.get('default')), "primary_key": c['name'] in pk_cols } for c in cols]
                    foreign_keys_info = [{"constrained_columns": fk['constrained_columns'], "referred_table": fk['referred_table'], "referred_columns": fk['referred_columns'] } for fk in fks]
                    total_rows = db.session.execute(text(f"SELECT COUNT(*) FROM {table_name}")).scalar(); active_rows = None
                    if any(c['name'] == 'is_deleted' for c in cols):
                         try: active_rows = db.session.execute(text(f"SELECT COUNT(*) FROM {table_name} WHERE is_deleted = FALSE")).scalar()
                         except Exception as count_err: logger.warning(f"Could not get active count for {table_name}: {count_err}")
                    table_info = {"columns": columns_info, "primary_keys": pk_cols, "foreign_keys": foreign_keys_info, "total_rows": total_rows }
                    if active_rows is not None: table_info["active_rows"] = active_rows; table_info["deleted_rows"] = total_rows - active_rows
                    schema_data[table_name] = table_info
                except Exception as table_err: logger.error(f"Error fetching schema for table {table_name}: {table_err}"); schema_data[table_name] = {"error": f"Failed schema: {table_err}"}
            db_name = db.engine.url.database; db_version = "N/A"
            try:
                if db.engine.dialect.name == 'postgresql': db_version = db.session.execute(text("SELECT version()")).scalar()
                elif db.engine.dialect.name == 'mysql': db_version = db.session.execute(text("SELECT VERSION()")).scalar()
            except Exception as db_info_err: logger.warning(f"Could not retrieve DB version: {db_info_err}")
            model_mapping = {m_cfg['model'].__tablename__: m_name for m_name, m_cfg in ReportService.ENTITY_CONFIG.items() if m_cfg.get('model')}
            response_data = {"database_info": {"name": db_name, "version": db_version, "total_tables": len(table_names), "application_models": len(model_mapping)}, "model_mapping": model_mapping, "tables": schema_data }
            return response_data, None
        except Exception as e: logger.exception(f"Failed to retrieve DB schema: {e}"); return None, f"Error retrieving schema: {e}"
    
    @staticmethod
    def _resolve_attribute_and_joins(base_model: Type, field_path: str, query: Query, current_aliases: Dict[str, Any]) -> Tuple[Query, Optional[Any], Dict[str, Any]]:
        """
        Helper to resolve a potentially nested attribute path (like 'form_submission.submitted_at')
        and ensure necessary joins are added to the query using aliases. It modifies the query
        object directly if joins are needed.

        Args:
            base_model: The starting SQLAlchemy model class (e.g., AnswerSubmitted).
            field_path: The dot-separated attribute path string (e.g., "form_submission.submitted_at").
            query: The current SQLAlchemy Query object to potentially modify.
            current_aliases: A dictionary tracking aliases already used for join paths.
                             Key: join path string (e.g., "form_submission"), Value: alias object.

        Returns:
            Tuple containing:
                - The potentially modified Query object (with joins added).
                - The resolved SQLAlchemy attribute object (e.g., AliasedClass.attribute) or None.
                - The updated dictionary of current_aliases.
        """
        parts = field_path.split('.')
        current_model_or_alias = base_model # Start with the base model type
        current_path_key = "" # Tracks the path for aliasing keys

        logger.debug(f"Resolving path: '{field_path}' starting from {base_model.__name__}")

        for i, part in enumerate(parts):
            is_last_part = (i == len(parts) - 1)

            # Inspect the current model or alias we are traversing from
            # If current_model_or_alias is an AliasedClass, inspect its mapper
            if isinstance(current_model_or_alias, type): # It's a model class
                 current_mapper = sqla_inspect(current_model_or_alias)
                 current_class_for_attr = current_model_or_alias
            else: # It's an AliasedClass instance
                 current_mapper = sqla_inspect(current_model_or_alias.entity)
                 current_class_for_attr = current_model_or_alias # Use alias directly for getattr

            if is_last_part:
                # Last part should be a column/attribute on the current model/alias
                if hasattr(current_class_for_attr, part):
                    resolved_attribute = getattr(current_class_for_attr, part)
                    logger.debug(f"Resolved attribute '{field_path}' to {resolved_attribute}")
                    # Return the query (possibly modified with joins) and the final attribute
                    return query, resolved_attribute, current_aliases
                else:
                    logger.warning(f"Attribute '{part}' not found on {current_mapper.class_.__name__} (aliased: {not isinstance(current_model_or_alias, type)}) for path '{field_path}'")
                    return query, None, current_aliases
            else:
                # Intermediate part must be a relationship defined on the current mapper
                if part in current_mapper.relationships:
                    relationship = current_mapper.relationships[part]
                    related_model = relationship.mapper.class_
                    relationship_key = relationship.key # The attribute name for the relationship (e.g., 'form_submission')

                    # Build the key for the alias dictionary (e.g., "form_submission" or "form_submission.form")
                    new_path_key = f"{current_path_key}.{relationship_key}" if current_path_key else relationship_key

                    if new_path_key not in current_aliases:
                        # Need to add a join with a new alias
                        logger.debug(f"Adding join for path '{new_path_key}' from {current_class_for_attr} using relationship '{relationship_key}'")
                        related_alias = aliased(related_model)
                        current_aliases[new_path_key] = related_alias # Store the new alias

                        # Join from the previous model/alias to the new alias using the relationship attribute
                        # Ensure we use the correct source for the relationship attribute
                        join_source = current_model_or_alias if not isinstance(current_model_or_alias, type) else base_model
                        query = query.join(related_alias, getattr(join_source, relationship_key))
                        current_model_or_alias = related_alias # Continue traversal from the new alias
                    else:
                        # Join already exists, just continue traversal from the existing alias
                        logger.debug(f"Reusing join alias for path '{new_path_key}'")
                        current_model_or_alias = current_aliases[new_path_key] # Continue from the existing alias

                    current_path_key = new_path_key # Update the path key for the next level

                else:
                    logger.warning(f"Relationship '{part}' not found on {current_mapper.class_.__name__} (aliased: {not isinstance(current_model_or_alias, type)}) for path '{field_path}'")
                    return query, None, current_aliases

        # Should not be reached if path is valid
        logger.error(f"Path resolution failed unexpectedly for '{field_path}'")
        return query, None, current_aliases